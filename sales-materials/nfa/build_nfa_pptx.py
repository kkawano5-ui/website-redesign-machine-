# -*- coding: utf-8 -*-
"""株式会社エヌエフエー様 AI活用・業務自動化 提案資料 PPTX"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY=RGBColor(0x13,0x29,0x4b); BLUE=RGBColor(0x1f,0x5f,0xa8); TEAL=RGBColor(0x1a,0xa3,0xa3)
BG=RGBColor(0xF4,0xF6,0xFA); CARD=RGBColor(0xFF,0xFF,0xFF); INK=RGBColor(0x1f,0x27,0x33)
SOFT=RGBColor(0x5b,0x66,0x75); LINE=RGBColor(0xE2,0xE7,0xEF); CHIP=RGBColor(0xE9,0xF1,0xFA)
RED=RGBColor(0xD6,0x45,0x45); AMBER=RGBColor(0xD9,0x92,0x11); GREEN=RGBColor(0x2f,0x9e,0x54)
WHITE=RGBColor(0xFF,0xFF,0xFF)
JP="Meiryo"
EW,EH=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=EW; prs.slide_height=EH
BLANK=prs.slide_layouts[6]

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EW,EH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s
def _round(sp,v=0.08):
    try: sp.adjustments[0]=v
    except Exception: pass
def rect(s,l,t,w,h,fill=None,line=None,lw=1.0,shp=MSO_SHAPE.RECTANGLE,rnd=0.0):
    sp=s.shapes.add_shape(shp,l,t,w,h)
    if rnd:_round(sp,rnd)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def text(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=4,ls=1.0):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(2); tf.margin_top=tf.margin_bottom=Pt(2)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.line_spacing=ls
        for (txt,sz,col,bold) in para:
            r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name=JP; rPr=r.font._rPr
            for tag in ("a:latin","a:ea","a:cs"):
                el=rPr.find(qn(tag))
                if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
                el.set("typeface",JP)
    return tb
def P(*r): return list(r)
def R(t,s,c=INK,b=False): return (t,s,c,b)
def header(s,eyebrow,heading,page):
    rect(s,Inches(0.7),Inches(0.6),Inches(0.16),Inches(0.5),fill=BLUE)
    text(s,Inches(0.98),Inches(0.53),Inches(11),Inches(0.35),[P(R(eyebrow,12,BLUE,True))])
    text(s,Inches(0.95),Inches(0.84),Inches(11.7),Inches(0.7),[P(R(heading,26,NAVY,True))])
    rect(s,Inches(0.98),Inches(1.48),Inches(2.0),Inches(0.035),fill=TEAL)
    text(s,Inches(12.4),Inches(7.04),Inches(0.7),Inches(0.3),[P(R(str(page),11,SOFT))],align=PP_ALIGN.RIGHT)
    text(s,Inches(0.7),Inches(7.04),Inches(7),Inches(0.3),[P(R("株式会社エヌエフエー様 ご提案",9,SOFT))])
def chip(s,l,t,txt,fill=BLUE,fg=WHITE):
    w=Inches(0.25+0.135*len(txt))
    rect(s,l,t,w,Inches(0.4),fill=fill,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.5)
    text(s,l,t,w,Inches(0.4),[P(R(txt,11.5,fg,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    return w
def card(s,l,t,w,h,title,body,icon=None,tcol=NAVY,fill=CARD,tsize=15,accent=None):
    rect(s,l,t,w,h,fill=fill,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.07)
    if accent: rect(s,l,t,w,Inches(0.12),fill=accent,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.4)
    yy=t.inches+0.2
    if icon:
        text(s,Inches(l.inches+0.25),Inches(yy),Inches(1),Inches(0.5),[P(R(icon,24,INK))],sa=0); yy+=0.5
    text(s,Inches(l.inches+0.25),Inches(yy),Inches(w.inches-0.45),Inches(0.6),[P(R(title,tsize,tcol,True))],sa=0,ls=1.05)
    text(s,Inches(l.inches+0.25),Inches(yy+(0.42 if icon else 0.5)),Inches(w.inches-0.45),h,[P(R(body,11.5,SOFT))],sa=0,ls=1.18)
L=Inches(0.7)

# ============ 1 COVER ============
s=slide(NAVY)
rect(s,0,0,EW,EH,fill=NAVY)
rect(s,0,Inches(4.7),EW,Inches(0.1),fill=TEAL)
rect(s,Inches(0.7),Inches(0.8),Inches(1.5),Inches(0.5),fill=None,line=TEAL,lw=1.5,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.3)
text(s,Inches(0.7),Inches(0.8),Inches(1.5),Inches(0.5),[P(R("AI × 自動化",12,TEAL,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
text(s,L,Inches(2.2),Inches(12),Inches(1.8),
     [P(R("AIで“全業務の自動化”を、",38,WHITE,True)),P(R("現場の標準装備に。",38,WHITE,True))],sa=2,ls=1.2)
text(s,L,Inches(4.0),Inches(12),Inches(0.5),[P(R("業務自動化・AI活用 ご提案 ／ AI研修のご案内",18,RGBColor(0xB8,0xC6,0xD9),False))],sa=0)
text(s,L,Inches(5.1),Inches(12),Inches(0.5),[P(R("株式会社エヌエフエー　代表取締役　大崎 玄長 様",17,WHITE,True))],sa=0)
text(s,L,Inches(6.5),Inches(12),Inches(0.4),[P(R("2026年 ／ ご提案・デモ資料",12,RGBColor(0x9F,0xB0,0xC6)))],sa=0)

# ============ 2 アジェンダ ============
s=slide()
header(s,"AGENDA ｜ 本日お話しすること","本日のゴールと進め方",2)
items=[("01","御社オペレーションの想定整理","30名で1,000名を支える体制を、業務マップで可視化"),
       ("02","自動化できる業務の棚卸し","“どこから手をつけるか”を効果×難易度で優先度づけ"),
       ("03","動くデモを2つご覧いただく","① 請求・給与ダブルチェック　② 求人原稿の自動生成"),
       ("04","AI研修での内製化プラン","外注ではなく、社員30名が“自分で作れる”状態へ")]
for i,(n,t,b) in enumerate(items):
    y=Inches(2.0+i*1.15)
    rect(s,L,y,Inches(11.9),Inches(1.0),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
    rect(s,Inches(0.9),Inches(y.inches+0.2),Inches(0.6),Inches(0.6),fill=NAVY,shp=MSO_SHAPE.OVAL)
    text(s,Inches(0.9),Inches(y.inches+0.2),Inches(0.6),Inches(0.6),[P(R(n,15,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(1.75),Inches(y.inches+0.15),Inches(10),Inches(0.45),[P(R(t,17,NAVY,True))],sa=0)
    text(s,Inches(1.75),Inches(y.inches+0.58),Inches(10),Inches(0.4),[P(R(b,12.5,SOFT))],sa=0)

# ============ 3 御社理解 ============
s=slide()
header(s,"UNDERSTANDING ｜ 御社の理解","私たちが理解している、御社の現在地",3)
text(s,L,Inches(1.95),Inches(7.0),Inches(2.4),
     [P(R("御社は大田区・蒲田を拠点に、人材派遣・紹介・教育に加え、研修センターや物流センターの直営運営まで手がける",13,SOFT),R("PEO型の総合人材サービス",13,NAVY,True),R("。",13,SOFT)),
      P(R("正社員30名というコンパクトな本部で、派遣社員約1,000名を支えています。すでに“全業務を自動化できないか”を社員に問い、請求書のダブルチェックなど具体的な取り組みも進行中——",13,SOFT)),
      P(R("AIを「一部の人の特技」から「30名全員の標準装備」へ広げるのが次の一手です。",13.5,NAVY,True))],sa=8,ls=1.45)
gx=Inches(8.1); gw=Inches(2.35)
for i,(v,u,lb) in enumerate([("30","名","正社員（本部）"),("1,000","名","派遣社員 規模"),("33","倍","1人が支える人数"),("PEO","型","総合人材サービス")]):
    col=i%2; row=i//2
    lx=Inches(8.1+col*2.55); ty=Inches(2.0+row*1.55)
    rect(s,lx,ty,gw,Inches(1.4),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.1)
    text(s,lx,Inches(ty.inches+0.2),gw,Inches(0.6),[P(R(v,30,BLUE,True),R(u,14,BLUE,True))],align=PP_ALIGN.CENTER,sa=0)
    text(s,lx,Inches(ty.inches+0.85),gw,Inches(0.4),[P(R(lb,11.5,SOFT))],align=PP_ALIGN.CENTER,sa=0)
rect(s,L,Inches(5.05),Inches(11.9),Inches(1.4),fill=CHIP,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
text(s,Inches(L.inches+0.3),Inches(5.22),Inches(11.3),Inches(1.1),
     [P(R("ポイント：",13,NAVY,True),R("少人数の本部が大人数の派遣を支える構造では、1人あたりの管理業務量が多く、",12.5,INK),R("“正確さ”と“スピード”の両立",12.5,NAVY,True),R("が常に求められます。ここが自動化の効果が最も大きい領域です。",12.5,INK))],sa=0,ls=1.4)

# ============ 4 構造的課題 ============
s=slide()
header(s,"ISSUE ｜ 派遣ビジネスの構造","“30名で1,000名”を支えるバックオフィスの負荷",4)
text(s,L,Inches(1.95),Inches(12),Inches(0.5),[P(R("派遣ビジネスは、派遣社員1人ごとに毎月くり返す定型業務が積み上がります。人数 × 毎月 × 正確性 が負荷の正体です。",13.5,SOFT))],sa=0)
loads=[("🕒","勤怠の収集・正規化","拠点・媒体ごとにバラバラな様式を毎月集約・転記"),
       ("🧾","請求・給与計算","単価×時間×手当を1,000名分。ミスが許されない"),
       ("✅","ダブルチェック","請求書・給与の突合を人手で二重確認"),
       ("📝","求人原稿・スカウト","媒体ごとに原稿作成、応募対応、面接調整"),
       ("📄","契約・コンプラ","派遣契約書、抵触日、36協定、社保手続き"),
       ("💬","問い合わせ対応","派遣社員からの給料日・有休・シフト等の質問対応")]
cw=Inches(3.78); ch=Inches(1.55)
for i,(ic,t,b) in enumerate(loads):
    col=i%3; row=i//3
    lx=Inches(0.7+col*4.05); ty=Inches(2.7+row*1.75)
    card(s,lx,ty,cw,ch,t,b,icon=None,tsize=14.5)
    text(s,Inches(lx.inches+0.25),Inches(ty.inches+0.18),Inches(1),Inches(0.4),[P(R(ic,20,INK))],sa=0)
    text(s,Inches(lx.inches+0.95),Inches(ty.inches+0.22),Inches(cw.inches-1.1),Inches(0.5),[P(R(t,14.5,NAVY,True))],sa=0)
    text(s,Inches(lx.inches+0.25),Inches(ty.inches+0.78),Inches(cw.inches-0.45),Inches(0.7),[P(R(b,11.5,SOFT))],sa=0,ls=1.18)

# ============ 5 業務マップ ============
s=slide()
header(s,"OPERATION MAP ｜ 業務全体像（想定）","御社オペレーションを5つの流れで想定",5)
flow=[("①","求人・採用","求人作成→応募対応→面接→登録・スキル管理",TEAL),
      ("②","マッチング・契約","ニーズと人材のマッチング→契約・抵触日管理",BLUE),
      ("③","勤怠・労務","勤怠収集→正規化→社保・36協定など労務管理",NAVY),
      ("④","請求・給与","勤怠×単価→請求/給与計算→ダブルチェック",BLUE),
      ("⑤","定着・教育","フォロー→問い合わせ対応→研修センター育成",TEAL)]
bw=Inches(2.28)
for i,(n,t,b,c) in enumerate(flow):
    lx=Inches(0.62+i*2.42); ty=Inches(2.3)
    rect(s,lx,ty,bw,Inches(2.7),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
    rect(s,lx,ty,bw,Inches(0.7),fill=c,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.18)
    text(s,lx,ty,bw,Inches(0.7),[P(R(n+"  "+t,14,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(lx.inches+0.2),Inches(ty.inches+0.95),Inches(bw.inches-0.4),Inches(1.6),[P(R(b,11.5,SOFT))],sa=0,ls=1.3)
    if i<4:
        text(s,Inches(lx.inches+bw.inches-0.05),Inches(ty.inches+1.0),Inches(0.45),Inches(0.5),[P(R("▶",16,BLUE,True))],align=PP_ALIGN.CENTER,sa=0)
rect(s,L,Inches(5.35),Inches(11.9),Inches(1.1),fill=CHIP,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
text(s,Inches(L.inches+0.3),Inches(5.52),Inches(11.3),Inches(0.8),
     [P(R("この5つの流れすべてに、AI・自動化で効く“すき間”があります。",13,NAVY,True),R("次ページで「効果 × 導入のしやすさ」で優先順位をつけます。",12.5,INK))],sa=0,ls=1.3)

# ============ 6 優先度マトリクス ============
s=slide()
header(s,"PRIORITY ｜ 自動化の機会マップ","“どこから始めるか” を一枚で",6)
ox,oy=Inches(2.6),Inches(2.05); ow,oh=Inches(8.6),Inches(4.6)
rect(s,ox,oy,ow,oh,fill=CARD,line=LINE)
# quadrant tint (top-right = best)
rect(s,Inches(ox.inches+ow.inches/2),oy,Inches(ow.inches/2),Inches(oh.inches/2),fill=RGBColor(0xEA,0xF6,0xEF))
# axes
rect(s,ox,Inches(oy.inches+oh.inches/2-0.01),ow,Inches(0.02),fill=LINE)
rect(s,Inches(ox.inches+ow.inches/2-0.01),oy,Inches(0.02),oh,fill=LINE)
text(s,Inches(1.0),oy,Inches(1.5),Inches(0.4),[P(R("効果：大 ▲",11,SOFT,True))],sa=0)
text(s,Inches(1.0),Inches(oy.inches+oh.inches-0.35),Inches(1.5),Inches(0.4),[P(R("効果：小 ▼",11,SOFT,True))],sa=0)
text(s,Inches(ox.inches+0.1),Inches(oy.inches+oh.inches+0.05),Inches(3),Inches(0.3),[P(R("◀ 導入：難しい",11,SOFT,True))],sa=0)
text(s,Inches(ox.inches+ow.inches-2.6),Inches(oy.inches+oh.inches+0.05),Inches(2.6),Inches(0.3),[P(R("導入：やさしい ▶",11,SOFT,True))],align=PP_ALIGN.RIGHT,sa=0)
text(s,Inches(ox.inches+ow.inches/2+0.1),Inches(oy.inches+0.08),Inches(4),Inches(0.3),[P(R("★ 最優先ゾーン（すぐ着手）",11,GREEN,True))],sa=0)
def bub(cx,cy,txt,c,w=2.0):
    lx=Inches(cx-w/2); ty=Inches(cy-0.28)
    rect(s,lx,ty,Inches(w),Inches(0.56),fill=c,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.5)
    text(s,lx,ty,Inches(w),Inches(0.56),[P(R(txt,10.5,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
# top-right (priority)
bub(9.4,2.9,"請求・給与チェック①",NAVY,2.3)
bub(9.0,3.7,"求人原稿生成②",TEAL,2.0)
bub(7.6,3.25,"問い合わせFAQ",BLUE,1.8)
# top-left (high effort high effect)
bub(4.3,2.85,"勤怠データ正規化",AMBER,2.1)
bub(4.0,3.75,"マッチング高度化",AMBER,2.1)
# bottom-right (easy, medium effect)
bub(9.2,5.5,"議事録・日報要約",BLUE,2.0)
bub(7.7,6.0,"メール下書き",BLUE,1.7)
# bottom-left
bub(4.2,5.7,"契約・抵触日管理",SOFT,2.1)
text(s,Inches(0.7),Inches(2.05),Inches(1.7),Inches(4.6),[P(R("①②は本日デモ",11,RED,True))],sa=0)

# ============ 7 テーマ1 請求チェック ============
s=slide()
header(s,"THEME 01 ｜ 本日デモ①","請求・給与のダブルチェックを自動化",7)
text(s,L,Inches(1.95),Inches(6.7),Inches(0.9),
     [P(R("社長が挙げられた“請求書のダブルチェック”そのものを自動化。勤怠×契約単価から正しい金額を自動計算し、人の確認が要る差異だけを浮かび上がらせます。",13,SOFT))],sa=0,ls=1.4)
for i,(t,b) in enumerate([("やること","勤怠データと契約マスタを照合し、請求額・給与額・粗利を自動算出"),
                          ("自動で検出","単価マスタ欠落／残業計算ミス／既存請求との差異／重複入力／労務リスク"),
                          ("人がやること","🔴と🟡で出た例外だけを確認。全件目視チェックが不要に"),
                          ("効果","チェック時間を大幅圧縮＋ヒューマンエラー由来の請求差異を防止")]):
    y=Inches(3.0+i*0.92)
    rect(s,L,y,Inches(2.0),Inches(0.78),fill=NAVY,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.12)
    text(s,L,y,Inches(2.0),Inches(0.78),[P(R(t,12.5,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    rect(s,Inches(2.85),y,Inches(4.55),Inches(0.78),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.1)
    text(s,Inches(3.05),y,Inches(4.2),Inches(0.78),[P(R(b,11.5,INK))],anchor=MSO_ANCHOR.MIDDLE,sa=0,ls=1.15)
# mock panel
mx=Inches(7.8)
rect(s,mx,Inches(2.0),Inches(4.8),Inches(4.5),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.04)
rect(s,mx,Inches(2.0),Inches(4.8),Inches(0.6),fill=NAVY,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.1)
text(s,mx,Inches(2.0),Inches(4.8),Inches(0.6),[P(R("自動チェックレポート（イメージ）",12.5,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
mock=[("🟢","佐藤 健","¥385,000","指摘なし",GREEN),("🔴","鈴木 一郎","¥462,800","既存請求と+¥17,200差異",RED),
      ("🟡","田中 翔","¥540,000","残業52h：労務リスク",AMBER),("🔴","伊藤 彩","—","契約マスタに単価なし",RED),
      ("🟡","高橋 美咲","¥319,000","勤怠の重複入力",AMBER)]
for i,(ic,nm,amt,fl,c) in enumerate(mock):
    y=Inches(2.8+i*0.7)
    rect(s,Inches(mx.inches+0.2),y,Inches(4.4),Inches(0.6),fill=BG,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.12)
    text(s,Inches(mx.inches+0.32),y,Inches(0.4),Inches(0.6),[P(R(ic,13,c))],anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(mx.inches+0.72),y,Inches(1.3),Inches(0.6),[P(R(nm,11.5,INK,True))],anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(mx.inches+1.95),y,Inches(1.0),Inches(0.6),[P(R(amt,11,INK))],anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(mx.inches+2.9),y,Inches(1.65),Inches(0.6),[P(R(fl,9.5,c,True))],anchor=MSO_ANCHOR.MIDDLE,sa=0,ls=1.0)

# ============ 8 テーマ2 求人原稿 ============
s=slide()
header(s,"THEME 02 ｜ 本日デモ②","求人原稿・スカウトをAIで一括生成",8)
text(s,L,Inches(1.95),Inches(12),Inches(0.6),
     [P(R("職種・時給・特徴を入れるだけで、媒体原稿・スカウトメール・SNS投稿を同時生成。採用スピードと応募数の両方に効きます。",13.5,SOFT))],sa=0)
for i,(t,b,c) in enumerate([("入力","職種／勤務地／時給／ターゲット／特徴を選ぶだけ",TEAL),
                            ("出力①","求人媒体の原稿（構成・キャッチ込み）",BLUE),
                            ("出力②","候補者向けスカウトメール文面",BLUE),
                            ("出力③","SNS投稿（ハッシュタグ付き）",BLUE)]):
    lx=Inches(0.7+i*3.05)
    card(s,lx,Inches(2.75),Inches(2.85),Inches(2.0),t,b,tsize=15,accent=c)
text(s,Inches(0.9),Inches(3.15),Inches(2.5),Inches(0.4),[P(R("",10))],sa=0)
rect(s,L,Inches(5.1),Inches(11.9),Inches(1.35),fill=CHIP,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
text(s,Inches(L.inches+0.3),Inches(5.27),Inches(11.3),Inches(1.05),
     [P(R("効果：",13,NAVY,True),R("1媒体30分→数分に短縮。応募実績の高い表現を学習させれば“勝ちパターン”を全担当で再現できます。担当者ごとの原稿品質のバラつきも解消。",12.5,INK))],sa=0,ls=1.4)
text(s,L,Inches(6.6),Inches(12),Inches(0.3),[P(R("※ 本番はClaude API連携で自然文を生成。デモはテンプレート方式で挙動をご覧いただきます。",9.5,SOFT))],sa=0)

# ============ 9 その他候補 ============
s=slide()
header(s,"MORE ｜ その他の自動化候補","まだまだ自動化できる業務があります",9)
more=[("💬","派遣社員 問い合わせ一次対応","給料日・有休・シフト等のFAQをチャットで自動回答。本部の電話/チャット対応を削減"),
      ("🕒","勤怠データの正規化","拠点・媒体ごとにバラバラな勤怠様式を、統一フォーマットへ自動変換・転記"),
      ("📄","契約書・抵触日チェック","派遣契約書の不備や抵触日の接近を自動アラート。コンプラ事故を未然防止"),
      ("🎓","研修教材・テスト生成","研修センター向けの教材・確認テスト・進捗レポートをAIで作成"),
      ("📝","議事録・日報の要約","商談メモや日報を自動要約。共有・引き継ぎの手間を削減"),
      ("📧","定型メール下書き","派遣先・候補者への連絡メールを状況に応じて自動下書き")]
cw=Inches(3.78); ch=Inches(1.75)
for i,(ic,t,b) in enumerate(more):
    col=i%3; row=i//3
    lx=Inches(0.7+col*4.05); ty=Inches(2.0+row*1.95)
    card(s,lx,ty,cw,ch,t,b,icon=ic,tsize=13.5)

# ============ 10 効果試算 ============
s=slide()
header(s,"IMPACT ｜ 想定効果","“塵も積もれば”を、数字で見る",10)
text(s,L,Inches(1.9),Inches(12),Inches(0.5),[P(R("以下は議論用のラフ試算です。実際の業務量を伺って精緻化します。",12.5,SOFT))],sa=0)
for i,(v,u,lb,c) in enumerate([("月100h+","","削減できる定型業務の目安*",BLUE),
                               ("0件","へ","請求の計算ミス・差異を抑制",GREEN),
                               ("数時間→数分","","求人原稿1本あたりの作成時間",TEAL),
                               ("30名","全員","“作れる人”を限定しない",NAVY)]):
    lx=Inches(0.7+i*3.05)
    rect(s,lx,Inches(2.5),Inches(2.85),Inches(1.9),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.08)
    rect(s,lx,Inches(2.5),Inches(2.85),Inches(0.12),fill=c,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.4)
    text(s,lx,Inches(2.95),Inches(2.85),Inches(0.7),[P(R(v,21,c,True),R(u,13,c,True))],align=PP_ALIGN.CENTER,sa=0)
    text(s,Inches(lx.inches+0.2),Inches(3.7),Inches(2.45),Inches(0.6),[P(R(lb,11.5,SOFT))],align=PP_ALIGN.CENTER,sa=0,ls=1.15)
rect(s,L,Inches(4.9),Inches(11.9),Inches(1.5),fill=CHIP,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
text(s,Inches(L.inches+0.3),Inches(5.08),Inches(11.3),Inches(1.2),
     [P(R("本質的な効果：",13,NAVY,True),R("削減した時間を“人にしかできない仕事”——派遣社員のフォロー、派遣先との関係構築、定着支援——に再投資できます。",12.5,INK)),
      P(R("自動化はコスト削減ではなく、御社の強み（人の手厚さ）を伸ばすための投資です。",12.5,NAVY,True))],sa=6,ls=1.35)
text(s,L,Inches(6.6),Inches(12),Inches(0.3),[P(R("* 1件6分×件数などの仮定に基づく目安。正式には御社実数で試算します。",9.5,SOFT))],sa=0)

# ============ 11 ロードマップ ============
s=slide()
header(s,"ROADMAP ｜ 進め方","AI研修 × 内製化｜90日の立ち上げ",11)
steps=[("STEP 1","〜30日","現状把握 & クイックウィン","業務の棚卸しワークショップ＋全社員AI基礎研修。すぐ効く1〜2業務を選定",TEAL),
       ("STEP 2","〜60日","デモ業務を実装・定着","請求チェック/求人生成を実運用へ。“作れる社員”を数名育成",BLUE),
       ("STEP 3","〜90日","横展開 & 内製文化","成功事例を全社共有。社員が自分で次の自動化を企画・作成できる状態へ",NAVY)]
for i,(st,d,t,b,c) in enumerate(steps):
    lx=Inches(0.7+i*4.05)
    rect(s,lx,Inches(2.2),Inches(3.8),Inches(3.4),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
    rect(s,lx,Inches(2.2),Inches(3.8),Inches(0.9),fill=c,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.12)
    text(s,lx,Inches(2.32),Inches(3.8),Inches(0.4),[P(R(st,15,WHITE,True))],align=PP_ALIGN.CENTER,sa=0)
    text(s,lx,Inches(2.72),Inches(3.8),Inches(0.35),[P(R(d,12,RGBColor(0xDD,0xE8,0xF2),True))],align=PP_ALIGN.CENTER,sa=0)
    text(s,Inches(lx.inches+0.3),Inches(3.35),Inches(3.2),Inches(0.6),[P(R(t,15,NAVY,True))],sa=0,ls=1.1)
    text(s,Inches(lx.inches+0.3),Inches(4.1),Inches(3.2),Inches(1.4),[P(R(b,12,SOFT))],sa=0,ls=1.3)
    if i<2: text(s,Inches(lx.inches+3.72),Inches(3.6),Inches(0.4),Inches(0.5),[P(R("▶",18,BLUE,True))],align=PP_ALIGN.CENTER,sa=0)
text(s,L,Inches(5.95),Inches(12),Inches(0.5),
     [P(R("ゴールは“ツール導入”ではなく、社員が自走してAIを使いこなす組織。だから外注ではなく研修なのです。",13,NAVY,True))],sa=0)

# ============ 12 なぜ研修 ============
s=slide()
header(s,"WHY TRAINING ｜ なぜ内製・研修なのか","“外注で1個作る”より“30名が作れる”",12)
rows=[("","外注・受託開発","AI研修で内製化"),
      ("作れる範囲","頼んだ1業務だけ","社員が気づいた業務すべて"),
      ("スピード","要件定義→納品で数ヶ月","その日のうちに試作・改善"),
      ("コスト","業務ごとに費用が発生","研修後は社内で内製＝低コスト"),
      ("定着","作った人がいないと変えられない","現場が自分で運用・改良"),
      ("御社らしさ","汎用的な仕組み","御社の業務・言葉に最適化")]
ty=2.0; rh=0.72
for i,(k,a,b) in enumerate(rows):
    y=Inches(ty+i*rh); head=(i==0)
    rect(s,L,y,Inches(2.7),Inches(rh-0.07),fill=(NAVY if head else CHIP),line=LINE)
    rect(s,Inches(L.inches+2.7),y,Inches(4.0),Inches(rh-0.07),fill=(NAVY if head else CARD),line=LINE)
    rect(s,Inches(L.inches+6.7),y,Inches(5.2),Inches(rh-0.07),fill=(NAVY if head else RGBColor(0xEA,0xF6,0xEF)),line=LINE)
    text(s,Inches(L.inches+0.15),y,Inches(2.5),Inches(rh),[P(R(k,12,(WHITE if head else NAVY),True))],anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(L.inches+2.85),y,Inches(3.7),Inches(rh),[P(R(a,12,(WHITE if head else SOFT),head))],anchor=MSO_ANCHOR.MIDDLE,sa=0,ls=1.1)
    text(s,Inches(L.inches+6.85),y,Inches(4.9),Inches(rh),[P(R(b,12,(WHITE if head else INK),True))],anchor=MSO_ANCHOR.MIDDLE,sa=0,ls=1.1)

# ============ 13 デモのご案内 ============
s=slide(NAVY)
rect(s,0,0,EW,EH,fill=NAVY)
text(s,L,Inches(0.8),Inches(11),Inches(0.4),[P(R("LIVE DEMO ｜ 本日ご覧いただくもの",13,TEAL,True))],sa=0)
text(s,L,Inches(1.35),Inches(12),Inches(0.8),[P(R("実際に、動かします。",32,WHITE,True))],sa=0)
demos=[("①","請求・給与 ダブルチェック自動化","勤怠×契約マスタ→請求/給与/粗利を自動計算し、差異・労務リスク・重複を自動検出。社長の“請求書のダブルチェック”そのものを自動化。",TEAL),
       ("②","求人原稿 自動生成","条件を選ぶだけで、媒体原稿・スカウト・SNS投稿を一括生成。採用のスピードと品質を底上げ。",BLUE)]
for i,(n,t,b,c) in enumerate(demos):
    y=Inches(2.6+i*1.9)
    rect(s,L,y,Inches(11.9),Inches(1.6),fill=RGBColor(0x1c,0x35,0x5e),shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
    rect(s,Inches(0.95),Inches(y.inches+0.35),Inches(0.9),Inches(0.9),fill=c,shp=MSO_SHAPE.OVAL)
    text(s,Inches(0.95),Inches(y.inches+0.35),Inches(0.9),Inches(0.9),[P(R(n,26,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(2.1),Inches(y.inches+0.25),Inches(10),Inches(0.5),[P(R(t,18,WHITE,True))],sa=0)
    text(s,Inches(2.1),Inches(y.inches+0.78),Inches(9.6),Inches(0.7),[P(R(b,12.5,RGBColor(0xC6,0xD5,0xE6)))],sa=0,ls=1.25)
text(s,L,Inches(6.5),Inches(12),Inches(0.4),[P(R("どちらもブラウザだけで動作。御社のデータ・職種に置き換えてその場で試せます。",12,TEAL,True))],sa=0)

# ============ 14 次のアクション ============
s=slide()
header(s,"NEXT ｜ 次のアクション","まずは“1業務”から、一緒に。",14)
for i,(t,b) in enumerate([("① 業務の棚卸しワークショップ","現場の声を集め、効果の大きい業務を一緒に特定（半日〜1日）"),
                          ("② クイックウィンを1つ実装","請求チェック等、すぐ効く業務をパイロット導入"),
                          ("③ 全社員向けAI研修をスタート","30名が“自分で作れる”状態へ。社内に推進役を育成")]):
    y=Inches(2.1+i*1.0)
    rect(s,L,y,Inches(8.5),Inches(0.85),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.1)
    text(s,Inches(L.inches+0.3),Inches(y.inches+0.12),Inches(8),Inches(0.4),[P(R(t,15,NAVY,True))],sa=0)
    text(s,Inches(L.inches+0.3),Inches(y.inches+0.5),Inches(8),Inches(0.35),[P(R(b,11.5,SOFT))],sa=0)
rect(s,Inches(9.5),Inches(2.1),Inches(3.1),Inches(2.75),fill=NAVY,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
text(s,Inches(9.5),Inches(2.5),Inches(3.1),Inches(0.5),[P(R("ご一緒できること",13,TEAL,True))],align=PP_ALIGN.CENTER,sa=0)
text(s,Inches(9.7),Inches(3.1),Inches(2.7),Inches(1.6),
     [P(R("✓ 業務棚卸し支援",12,WHITE,True)),P(R("✓ AI研修の設計・実施",12,WHITE,True)),P(R("✓ 自動化ツールの内製伴走",12,WHITE,True))],sa=8,ls=1.2)
rect(s,L,Inches(5.4),Inches(11.9),Inches(1.1),fill=CHIP,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
text(s,Inches(L.inches+0.3),Inches(5.62),Inches(11.3),Inches(0.7),
     [P(R("“全業務を自動化できないか常に意識する”という御社の文化に、AIという武器を。",13.5,NAVY,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)

prs.save("/home/user/website-redesign-machine-/sales-materials/nfa/NFA様_AI業務自動化_提案資料.pptx")
print("slides:",len(prs.slides._sldIdLst))
