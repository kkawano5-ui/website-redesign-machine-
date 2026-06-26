# -*- coding: utf-8 -*-
"""陶彩館 AI研修デモ用 営業資料 PPTX ジェネレーター"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
TERRA   = RGBColor(0xB5,0x53,0x2F)
TERRA_D = RGBColor(0x8C,0x3C,0x1F)
AMBER   = RGBColor(0xD9,0x8A,0x3D)
CREAM   = RGBColor(0xFA,0xF5,0xEE)
CARD    = RGBColor(0xFF,0xFF,0xFF)
INK     = RGBColor(0x3A,0x32,0x2C)
INKSOFT = RGBColor(0x6B,0x5F,0x54)
LINE    = RGBColor(0xE6,0xDC,0xCD)
GREEN   = RGBColor(0x5B,0x7A,0x4F)
CHIPBG  = RGBColor(0xFB,0xEE,0xDE)
WHITE   = RGBColor(0xFF,0xFF,0xFF)

JP = "Meiryo"
EW, EH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = EW
prs.slide_height = EH
BLANK = prs.slide_layouts[6]

def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EW,EH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def _set_round(shape, val=0.08):
    try: shape.adjustments[0] = val
    except Exception: pass

def rect(s,l,t,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,round_=0.0,shadow=False):
    sp = s.shapes.add_shape(shape,l,t,w,h)
    if round_: _set_round(sp, round_)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp

def text(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,
         space_after=4,line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (txt,size,color,bold)"""
    tb = s.shapes.add_textbox(l,t,w,h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=tf.margin_right=Pt(2); tf.margin_top=tf.margin_bottom=Pt(2)
    for i,para in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt,size,color,bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = JP
            rPr = r.font._rPr
            for tag in ("a:latin","a:ea","a:cs"):
                el = rPr.find(qn(tag))
                if el is None:
                    el = rPr.makeelement(qn(tag), {}); rPr.append(el)
                el.set("typeface", JP)
    return tb

def P(*runs): return list(runs)         # paragraph builder
def R(txt,size,color=INK,bold=False): return (txt,size,color,bold)

def header(s, eyebrow, heading, page):
    rect(s, Inches(0.7), Inches(0.62), Inches(0.16), Inches(0.5), fill=TERRA)
    text(s, Inches(0.98), Inches(0.55), Inches(11), Inches(0.35),
         [P(R(eyebrow, 12, TERRA, True))])
    text(s, Inches(0.95), Inches(0.86), Inches(11.6), Inches(0.7),
         [P(R(heading, 27, INK, True))])
    rect(s, Inches(0.98), Inches(1.52), Inches(2.2), Inches(0.035), fill=TERRA)
    text(s, Inches(12.5), Inches(7.02), Inches(0.6), Inches(0.3),
         [P(R(str(page), 11, INKSOFT))], align=PP_ALIGN.RIGHT)
    text(s, Inches(0.7), Inches(7.02), Inches(5), Inches(0.3),
         [P(R("陶彩館｜オーヤシマ株式会社", 9, INKSOFT))])

def chip(s,l,t,txt,fill=TERRA,fg=WHITE,w=None):
    w = w or Inches(0.2+0.14*len(txt))
    c = rect(s,l,t,w,Inches(0.42),fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.5)
    text(s,l,t,w,Inches(0.42),[P(R(txt,12,fg,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=0)
    return w

def stat(s,l,t,w,num,unit,lbl):
    rect(s,l,t,w,Inches(1.35),fill=CARD,line=LINE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.12)
    text(s,l,Inches(t.inches+0.18),w,Inches(0.6),
         [P(R(num,33,TERRA,True),R(unit,15,TERRA,True))],align=PP_ALIGN.CENTER,space_after=0)
    text(s,l,Inches(t.inches+0.82),w,Inches(0.5),
         [P(R(lbl,11.5,INKSOFT))],align=PP_ALIGN.CENTER,space_after=0,line_spacing=1.0)

def card(s,l,t,w,h,title,body,icon=None,tcolor=TERRA_D,fill=CARD,line_=LINE,tsize=15):
    rect(s,l,t,w,h,fill=fill,line=line_,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.08)
    yy = t.inches+0.22
    if icon:
        text(s,l,Inches(yy),w,Inches(0.5),[P(R(icon,26,INK))],align=PP_ALIGN.CENTER if False else PP_ALIGN.LEFT,space_after=0)
        yy += 0.55
    text(s,Inches(l.inches+0.05),Inches(yy),Inches(w.inches-0.4),Inches(0.6),
         [P(R(title,tsize,tcolor,True))],space_after=0,line_spacing=1.05)
    text(s,Inches(l.inches+0.05),Inches(yy+0.5 if not icon else yy+0.42),Inches(w.inches-0.3),h,
         [P(R(body,11.5,INKSOFT))],space_after=0,line_spacing=1.15)

L = Inches(0.7)  # left margin

# ===================== 1. COVER =====================
s = slide(CREAM)
rect(s,0,0,EW,EH,fill=RGBColor(0xF6,0xE7,0xD6))
rect(s,0,0,EW,Inches(2.6),fill=TERRA)
rect(s,0,Inches(2.6),EW,Inches(0.12),fill=AMBER)
text(s,L,Inches(0.7),Inches(11),Inches(1.0),[P(R("陶彩館",46,WHITE,True))],space_after=0)
text(s,L,Inches(1.72),Inches(11),Inches(0.4),[P(R("T O U S A I K A N  ｜  和歌山・泉佐野のリフォーム専門店",13,RGBColor(0xF3,0xE0,0xCF),True))],space_after=0)
text(s,L,Inches(3.15),Inches(12),Inches(1.2),[P(R("住まいと暮らしに、彩りを。",40,INK,True))],space_after=0)
text(s,L,Inches(4.25),Inches(12),Inches(0.6),[P(R("― 会社紹介・営業提案資料 ―",19,INKSOFT))],space_after=0)
x=L.inches
for txt in ["創業50年","施工実績 55,000件超","和歌山県 売上No.1"]:
    w=chip(s,Inches(x),Inches(5.2),txt); x+=w.inches+0.2
chip(s,Inches(x),Inches(5.2),"自社施工・安心保証",fill=WHITE,fg=TERRA)
text(s,L,Inches(6.55),Inches(11),Inches(0.4),[P(R("オーヤシマ株式会社",13,INKSOFT,True))],space_after=0)

# ===================== 2. 想い / 理念 =====================
s = slide()
header(s,"OUR MISSION ｜ 私たちの想い","家族の笑顔を、地域と未来へつなぐ",2)
text(s,L,Inches(2.0),Inches(6.6),Inches(0.4),[P(R("経営理念",14,INKSOFT,True))],space_after=0)
text(s,L,Inches(2.4),Inches(6.6),Inches(1.4),
     [P(R("「住まいと暮らしに彩りを",22,TERRA_D,True)),
      P(R("〜家族の笑顔を地域社会と未来につなぐ〜」",22,TERRA_D,True))],space_after=2,line_spacing=1.25)
text(s,L,Inches(4.1),Inches(6.6),Inches(2.6),
     [P(R("1975年の創業以来、半世紀にわたり地域の住まいと向き合ってきました。リフォームは「工事」ではなく、お客様一人ひとりの暮らしの想いをかたちにすること。水まわりの小さなお困りごとから、住み継ぐための大規模な改修まで、「相談してよかった」と言っていただける一社であり続けます。",13.5,INKSOFT))],
     space_after=0,line_spacing=1.4)
gx=Inches(8.0); gw=Inches(2.35)
stat(s,gx,Inches(2.1),gw,"50","年","地域の実績（1975年〜）")
stat(s,Inches(gx.inches+2.55),Inches(2.1),gw,"55,000","件+","累計施工実績")
stat(s,gx,Inches(3.7),gw,"No.1","","和歌山県 売上")
stat(s,Inches(gx.inches+2.55),Inches(3.7),gw,"2","店舗","和歌山・大阪に展開")

# ===================== 3. 会社概要 =====================
s = slide()
header(s,"COMPANY ｜ 会社概要","会社概要",3)
rows=[("商号","オーヤシマ株式会社（リフォーム専門店「陶彩館」）"),
      ("創業","1975年（昭和50年）　― 半世紀にわたり地域の住まいを支える"),
      ("事業内容","住宅リフォーム全般／水まわり設備工事／屋根・外壁塗装／増改築・断熱改修／バリアフリー改修／太陽光・省エネ設備／外構工事 ほか"),
      ("拠点","■ 岩出本店（ショールーム）和歌山県岩出市　　■ りんくう店 大阪府泉佐野市"),
      ("ショールーム","展示面積 約800坪・実機80台以上を常設展示（キッチン／バス／トイレ／洗面 ほか）"),
      ("体制","ショールームアドバイザー・設計・コーディネーター・施工監理・施工・メンテナンスの各専門部門による一貫体制"),
      ("加盟・パートナー","タカラスタンダード パートナーショップ／TOTO リモデルクラブ店／Panasonic リフォーム ほか")]
ty=2.0; rh=0.62
for i,(k,v) in enumerate(rows):
    yy=Inches(ty+i*rh)
    rect(s,L,yy,Inches(2.5),Inches(rh-0.06),fill=CHIPBG,line=LINE)
    rect(s,Inches(L.inches+2.5),yy,Inches(9.3),Inches(rh-0.06),fill=CARD,line=LINE)
    text(s,Inches(L.inches+0.12),yy,Inches(2.3),Inches(rh),[P(R(k,12.5,TERRA_D,True))],anchor=MSO_ANCHOR.MIDDLE,space_after=0)
    text(s,Inches(L.inches+2.65),yy,Inches(9.0),Inches(rh),[P(R(v,11.5,INK))],anchor=MSO_ANCHOR.MIDDLE,space_after=0,line_spacing=1.05)
text(s,L,Inches(6.5),Inches(12),Inches(0.4),
     [P(R("※ 本資料の数値・項目は公開情報をもとにした営業提案用サンプルです。正式利用時は社内の最新情報に差し替えてください。",9.5,INKSOFT))],space_after=0)

# ===================== 4. 選ばれる3つの理由 =====================
s = slide()
header(s,"WHY US ｜ 選ばれ続ける理由","陶彩館が選ばれる、3つの理由",4)
cw=Inches(3.75); cy=Inches(2.1); ch=Inches(3.7)
data=[("🏆","① 確かな実績","創業50年・施工実績55,000件超。和歌山県売上No.1の信頼。地域で積み重ねた経験が、最適なご提案を支えます。"),
      ("🏠","② 体感できる大型ショールーム","約800坪・実機80台以上。実物を見て・ふれて・比べて選べるから、仕上がりのイメージ違いがありません。"),
      ("🤝","③ 自社施工＆ワンストップ","相談・設計・施工・アフターまで自社一貫。工事を丸投げしないから、責任の所在が明確で安心です。")]
for i,(ic,t,b) in enumerate(data):
    lx=Inches(0.7+i*4.0)
    card(s,lx,cy,cw,ch,t,b,icon=ic,tsize=17)
text(s,L,Inches(6.05),Inches(12),Inches(0.4),
     [P(R("― 次ページから、3つの理由を詳しくご紹介します ―",15,INKSOFT))],space_after=0)

# ===================== 5. 理由1 実績 =====================
s = slide()
header(s,"REASON 01 ｜ 確かな実績","50年・55,000件が証明する信頼",5)
text(s,L,Inches(2.05),Inches(6.7),Inches(2.0),
     [P(R("リフォームは、暮らしてみて初めて分かることが多いお買い物です。だからこそ、地域で長く・数多く手がけてきた経験がものを言います。陶彩館は和歌山・泉佐野エリアで圧倒的な施工件数を積み重ね、「この間取りなら」「この築年数なら」という引き出しの多さでご提案します。",13.5,INKSOFT))],
     space_after=0,line_spacing=1.45)
ticks=["和歌山県 リフォーム売上 No.1の実績","水まわりから屋根・外壁、増改築まで幅広い対応力","地域密着だから、施工後も「すぐ駆けつけられる」距離感"]
for i,tk in enumerate(ticks):
    yy=Inches(4.1+i*0.55)
    rect(s,L,Inches(yy.inches+0.02),Inches(0.28),Inches(0.28),fill=GREEN,shape=MSO_SHAPE.OVAL)
    text(s,L,Inches(yy.inches-0.02),Inches(0.28),Inches(0.3),[P(R("✓",11,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=0)
    text(s,Inches(L.inches+0.42),yy,Inches(6.3),Inches(0.5),[P(R(tk,13,INK))],space_after=0)
gx=Inches(8.3); gw=Inches(2.2)
stat(s,gx,Inches(2.1),gw,"55,000","件+","累計施工実績")
stat(s,Inches(gx.inches+2.4),Inches(2.1),gw,"50","年","創業からの歩み")
stat(s,gx,Inches(3.7),gw,"No.1","","和歌山県 売上")
stat(s,Inches(gx.inches+2.4),Inches(3.7),gw,"14","億円+","年商規模")

# ===================== 6. 理由2 ショールーム =====================
s = slide()
header(s,"REASON 02 ｜ 体感できるショールーム","「見て・ふれて・比べて」選べる安心",6)
text(s,L,Inches(2.0),Inches(12),Inches(0.6),
     [P(R("カタログだけでは分からない、サイズ感・質感・使い勝手。約800坪の大型展示で“実物”をご体感いただけます。",14,INKSOFT))],space_after=0)
gw=Inches(2.78)
stat(s,Inches(0.7),Inches(2.7),gw,"約800","坪","展示面積")
stat(s,Inches(3.7),Inches(2.7),gw,"80","台+","実機展示")
stat(s,Inches(6.7),Inches(2.7),gw,"複数","","メーカーを一度に比較")
stat(s,Inches(9.7),Inches(2.7),gw,"予約","制","じっくり相談")
cw=Inches(2.78); cy=Inches(4.35); ch=Inches(2.1)
cd=[("キッチン","システムキッチンを多数展示。高さ・収納・動線を実感。"),
    ("バスルーム","システムバスの広さ・保温性・お掃除のしやすさを体感。"),
    ("トイレ・洗面","最新の節水・お手入れ性能を見比べて選べます。"),
    ("専門アドバイザー","各分野の担当が、暮らしに合う1台をご提案。")]
for i,(t,b) in enumerate(cd):
    card(s,Inches(0.7+i*3.0),cy,cw,ch,t,b,tsize=15)

# ===================== 7. 理由3 自社施工 =====================
s = slide()
header(s,"REASON 03 ｜ 自社施工・ワンストップ体制","相談から完成・その後まで、一社完結",7)
text(s,L,Inches(2.05),Inches(6.6),Inches(2.2),
     [P(R("陶彩館は工事を外注に丸投げせず、自社施工で品質を管理します。営業・設計・コーディネート・施工監理・施工・アフターまで、各分野の専門スタッフがチームでお客様を担当。「言った・言わない」「窓口がたらい回し」がなく、責任の所在が明確です。",13.5,INKSOFT))],
     space_after=0,line_spacing=1.45)
ticks=["窓口が一本化され、打合せがスムーズ","現場を知る自社職人だから仕上がりが安定","引き渡し後のメンテナンスも自社で対応"]
for i,tk in enumerate(ticks):
    yy=Inches(4.25+i*0.55)
    rect(s,L,Inches(yy.inches+0.02),Inches(0.28),Inches(0.28),fill=GREEN,shape=MSO_SHAPE.OVAL)
    text(s,L,Inches(yy.inches-0.02),Inches(0.28),Inches(0.3),[P(R("✓",11,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=0)
    text(s,Inches(L.inches+0.42),yy,Inches(6.1),Inches(0.5),[P(R(tk,13,INK))],space_after=0)
bx=Inches(7.9); bw=Inches(4.7)
rect(s,bx,Inches(2.1),bw,Inches(4.3),fill=CHIPBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.05)
text(s,bx,Inches(2.28),bw,Inches(0.5),[P(R("専門チームでお客様を担当",16,INK,True))],align=PP_ALIGN.CENTER,space_after=0)
team=["ショールーム／リフォームアドバイザー","設計・コーディネーター","施工監理・自社施工"]
for i,tm in enumerate(team):
    yy=Inches(2.95+i*0.78)
    rect(s,Inches(bx.inches+0.3),yy,Inches(bw.inches-0.6),Inches(0.62),fill=CARD,line=LINE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.12)
    text(s,Inches(bx.inches+0.3),yy,Inches(bw.inches-0.6),Inches(0.62),[P(R(tm,13,INK,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=0)
yy=Inches(5.35)
rect(s,Inches(bx.inches+0.3),yy,Inches(bw.inches-0.6),Inches(0.7),fill=GREEN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.12)
text(s,Inches(bx.inches+0.3),yy,Inches(bw.inches-0.6),Inches(0.7),[P(R("アフター・メンテナンス専任スタッフ",13,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=0)

# ===================== 8. 商材ラインナップ =====================
s = slide()
header(s,"PRODUCTS ｜ 取扱い商材","住まいのことなら、何でもご相談ください",8)
prod=[("🚰","水まわり","キッチン・浴室・トイレ・洗面。満足度を一番上げやすい工事。"),
      ("🎨","屋根・外壁塗装","住まいを雨風から守り外観も一新。建物の寿命を延ばします。"),
      ("🏗️","増改築・間取り変更","家族構成や暮らし方の変化に合わせ住まいを再設計。"),
      ("🌡️","断熱・省エネ改修","夏涼しく冬暖かい家へ。光熱費とからだにやさしく。"),
      ("♿","バリアフリー改修","手すり・段差解消。これからも安心して暮らせる住まいへ。"),
      ("☀️","太陽光・蓄電","創エネ・省エネで電気代と災害への備えを両立。"),
      ("🌿","外構・エクステリア","玄関まわり・カーポート・お庭まで第一印象を整える。"),
      ("🛠️","小工事・修繕","「ここだけ」のお困りごとも歓迎。お付き合いの入口に。")]
cw=Inches(2.85); ch=Inches(1.95)
for i,(ic,t,b) in enumerate(prod):
    col=i%4; row=i//4
    lx=Inches(0.7+col*3.07); ty=Inches(2.05+row*2.1)
    card(s,lx,ty,cw,ch,t,b,icon=ic,tsize=14.5)
text(s,L,Inches(6.45),Inches(12),Inches(0.4),
     [P(R("水まわりの1ヶ所から、住まい全体のトータルリフォームまで。“まとめて相談できる”のが強みです。",13,TERRA_D,True))],space_after=0)

# ===================== 9. 取扱メーカー =====================
s = slide()
header(s,"MAKERS ｜ 取扱メーカー","主要メーカー 正規パートナー",9)
text(s,L,Inches(2.0),Inches(12),Inches(0.6),
     [P(R("国内主要メーカーの認定店だから、豊富な品揃え・適正価格・メーカー保証まで安心。1社では選びきれない比較も、陶彩館なら一度に。",13.5,INKSOFT))],space_after=0)
mk=[("タカラスタンダード","ホーロー素材に強み。汚れに強く長持ちするキッチン・バス。","パートナーショップ"),
    ("TOTO","節水・清掃性に優れた水まわり設備。","リモデルクラブ店"),
    ("Panasonic","住宅設備をトータルに展開。","リフォームパートナー")]
cw=Inches(3.85); ch=Inches(2.0)
for i,(t,b,tag) in enumerate(mk):
    lx=Inches(0.7+i*4.05)
    rect(s,lx,Inches(2.75),cw,ch,fill=CARD,line=LINE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.08)
    text(s,lx,Inches(2.95),cw,Inches(0.5),[P(R(t,17,TERRA_D,True))],align=PP_ALIGN.CENTER,space_after=0)
    text(s,Inches(lx.inches+0.25),Inches(3.5),Inches(cw.inches-0.5),Inches(0.9),[P(R(b,12,INKSOFT))],align=PP_ALIGN.CENTER,space_after=0,line_spacing=1.2)
    pw=Inches(0.3+0.14*len(tag))
    chip(s,Inches(lx.inches+(cw.inches-pw.inches)/2),Inches(4.3),tag,fill=CHIPBG,fg=TERRA_D,w=pw)
rect(s,L,Inches(5.15),Inches(11.93),Inches(1.25),fill=CHIPBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.06)
text(s,Inches(L.inches+0.3),Inches(5.32),Inches(11.4),Inches(1.0),
     [P(R("営業トーク例：",12.5,TERRA_D,True),R("「メーカーごとに“得意分野”が違います。タカラさんは汚れに強いホーロー、TOTOさんは節水とお掃除のしやすさ…。当店なら実機を並べて比較できるので、ご予算とご希望に一番合う1台を一緒に選べますよ。」",12.5,INK))],
     space_after=0,line_spacing=1.25)
text(s,L,Inches(6.55),Inches(12),Inches(0.3),
     [P(R("※ 取扱メーカー・加盟内容は時期により異なります。最新状況は店舗にご確認ください。",9.5,INKSOFT))],space_after=0)

# ===================== 10. 流れ =====================
s = slide()
header(s,"FLOW ｜ ご相談から完成までの流れ","はじめてでも安心の6ステップ",10)
steps=[("1","ご相談・ショールーム来店","お困りごとやご希望をヒアリング。実機を見ながらイメージを膨らませます。"),
       ("2","現地調査・お見積り","専門スタッフが現場を確認し、最適なプランとお見積りをご提示。"),
       ("3","プランご提案・ご契約","設計・コーディネーターが仕様を具体化。ご納得いただいてからご契約。"),
       ("4","着工・自社施工","施工監理のもと自社職人が丁寧に工事。近隣への配慮も徹底。"),
       ("5","完成・お引き渡し","仕上がりをご一緒に確認。使い方もしっかりご説明します。"),
       ("6","アフター・メンテナンス","工事保証と専任スタッフで、引き渡し後も末永くサポート。")]
for i,(n,t,b) in enumerate(steps):
    col=i//3; row=i%3
    lx=Inches(0.8+col*6.1); ty=Inches(2.2+row*1.55)
    rect(s,lx,ty,Inches(0.62),Inches(0.62),fill=TERRA,shape=MSO_SHAPE.OVAL)
    text(s,lx,ty,Inches(0.62),Inches(0.62),[P(R(n,20,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=0)
    text(s,Inches(lx.inches+0.85),Inches(ty.inches-0.02),Inches(4.9),Inches(0.4),[P(R(t,15.5,INK,True))],space_after=0)
    text(s,Inches(lx.inches+0.85),Inches(ty.inches+0.42),Inches(4.9),Inches(0.9),[P(R(b,11.5,INKSOFT))],space_after=0,line_spacing=1.2)

# ===================== 11. 保証・アフター =====================
s = slide()
header(s,"WARRANTY ｜ 保証・アフターサービス","「つくって終わり」にしない安心",11)
ticks=["工事保証で、施工後の万一に備えます","メーカー保証＋正規パートナーならではの確実な対応","専任メンテナンススタッフがアフターを担当","地域密着だから、何かあればすぐ駆けつけられる","長いお付き合いの中で、次の住まいの相談先に"]
for i,tk in enumerate(ticks):
    yy=Inches(2.2+i*0.7)
    rect(s,L,Inches(yy.inches+0.03),Inches(0.3),Inches(0.3),fill=GREEN,shape=MSO_SHAPE.OVAL)
    text(s,L,Inches(yy.inches),Inches(0.3),Inches(0.3),[P(R("✓",12,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=0)
    text(s,Inches(L.inches+0.45),yy,Inches(6.0),Inches(0.6),[P(R(tk,14,INK))],space_after=0)
qx=Inches(7.5); qw=Inches(5.1)
for i,(q,who) in enumerate([("「水まわりのリフォームでお願いしました。打合せから施工まで同じ会社で、対応も丁寧。引き渡し後の点検まで安心でした。」","― お客様の声（イメージ）"),
                            ("「外注に丸投げしないと聞いて決めました。職人さんの仕事が丁寧で、近所への気配りも嬉しかったです。」","― お客様の声（イメージ）")]):
    ty=Inches(2.2+i*1.85)
    rect(s,qx,ty,Inches(0.1),Inches(1.55),fill=AMBER)
    rect(s,Inches(qx.inches+0.1),ty,qw,Inches(1.55),fill=CARD,line=LINE)
    text(s,Inches(qx.inches+0.35),Inches(ty.inches+0.15),Inches(qw.inches-0.5),Inches(1.0),[P(R(q,12.5,INK))],space_after=0,line_spacing=1.25)
    text(s,Inches(qx.inches+0.35),Inches(ty.inches+1.18),Inches(qw.inches-0.5),Inches(0.3),[P(R(who,10.5,INKSOFT))],space_after=0)
text(s,L,Inches(6.55),Inches(12),Inches(0.3),
     [P(R("※ お客様の声は提案イメージ用のサンプルです。掲載時は許諾を得た実例に差し替えてください。",9.5,INKSOFT))],space_after=0)

# ===================== 12. 施工事例 =====================
s = slide()
header(s,"CASES ｜ 施工事例","暮らしが変わった、リフォーム事例",12)
cases=[(RGBColor(0xD9,0xB4,0x8F),"キッチン全面リフォーム","古い独立キッチンを対面式に。家族との会話が増える明るいLDKへ。"),
       (RGBColor(0x9B,0xBD,0x9B),"浴室まるごと交換","寒く滑りやすい在来浴室を、暖かく安全なシステムバスへ。"),
       (RGBColor(0xCB,0xB0,0x6A),"屋根・外壁塗装","色あせた外観を一新し、雨漏りリスクも解消。資産価値を維持。")]
cw=Inches(3.85)
for i,(col,t,b) in enumerate(cases):
    lx=Inches(0.7+i*4.05)
    rect(s,lx,Inches(2.1),cw,Inches(2.9),fill=CARD,line=LINE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.06)
    rect(s,Inches(lx.inches+0.2),Inches(2.3),Inches(cw.inches-0.4),Inches(1.1),fill=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.08)
    text(s,Inches(lx.inches+0.2),Inches(2.3),Inches(cw.inches-0.4),Inches(1.1),[P(R("Before  →  After",14,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=0)
    text(s,Inches(lx.inches+0.25),Inches(3.55),Inches(cw.inches-0.5),Inches(0.5),[P(R(t,15.5,TERRA_D,True))],space_after=0)
    text(s,Inches(lx.inches+0.25),Inches(4.05),Inches(cw.inches-0.5),Inches(0.9),[P(R(b,11.5,INKSOFT))],space_after=0,line_spacing=1.2)
rect(s,L,Inches(5.25),Inches(11.93),Inches(1.15),fill=CHIPBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.06)
text(s,Inches(L.inches+0.3),Inches(5.42),Inches(11.4),Inches(0.9),
     [P(R("営業ポイント：",12.5,TERRA_D,True),R("提案時は「お客様の家と築年数・間取りが近い事例」を見せるのが効果的。施工事例ページ（tousaikan.com/case）から該当事例を選び、Before/Afterと費用感をセットで提示しましょう。",12.5,INK))],
     space_after=0,line_spacing=1.25)

# ===================== 13. 店舗・エリア =====================
s = slide()
header(s,"STORES ｜ 店舗・対応エリア","和歌山・大阪の2店舗で地域をカバー",13)
stores=[("🏠 岩出本店（ショールーム）","和歌山県岩出市","対応エリア：岩出市・和歌山市・紀の川市・橋本市・有田市 ほか和歌山県内","約800坪の大型ショールームを併設。実機を見比べながらじっくりご相談いただけます。"),
        ("🏠 りんくう店","大阪府泉佐野市","対応エリア：泉佐野市・阪南市・泉南市・貝塚市・泉南郡 ほか泉州エリア","大阪・泉州エリアのお客様の窓口。和歌山本店と連携してご対応します。")]
cw=Inches(5.85); ch=Inches(3.6)
for i,(t,loc,area,desc) in enumerate(stores):
    lx=Inches(0.7+i*6.13)
    rect(s,lx,Inches(2.1),cw,ch,fill=CARD,line=LINE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.05)
    text(s,Inches(lx.inches+0.35),Inches(2.35),Inches(cw.inches-0.6),Inches(0.5),[P(R(t,18,TERRA_D,True))],space_after=0)
    text(s,Inches(lx.inches+0.35),Inches(3.0),Inches(cw.inches-0.6),Inches(0.4),[P(R(loc,15,INK,True))],space_after=0)
    text(s,Inches(lx.inches+0.35),Inches(3.6),Inches(cw.inches-0.6),Inches(0.9),[P(R(area,12.5,INKSOFT))],space_after=0,line_spacing=1.25)
    text(s,Inches(lx.inches+0.35),Inches(4.6),Inches(cw.inches-0.6),Inches(0.9),[P(R(desc,12.5,INKSOFT))],space_after=0,line_spacing=1.25)
text(s,L,Inches(6.0),Inches(12),Inches(0.5),
     [P(R("ご来店が難しい場合は、現地調査・出張相談も承ります。まずはお気軽にお問い合わせください。",13.5,TERRA_D,True))],space_after=0)

# ===================== 14. CTA =====================
s = slide(CREAM)
rect(s,0,0,EW,EH,fill=RGBColor(0xF6,0xE7,0xD6))
text(s,L,Inches(0.7),Inches(11),Inches(0.4),[P(R("CONTACT ｜ お問い合わせ",13,TERRA,True))],space_after=0)
text(s,L,Inches(1.3),Inches(12),Inches(1.6),
     [P(R("まずは、ショールームへ。",34,INK,True)),
      P(R("「見て・ふれて・比べる」ことから始めましょう。",34,INK,True))],space_after=2,line_spacing=1.2)
rect(s,L,Inches(3.55),Inches(9.5),Inches(2.1),fill=TERRA,shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.06)
text(s,Inches(L.inches+0.5),Inches(3.8),Inches(8.5),Inches(0.5),[P(R("ご相談・お見積り・ショールーム来店予約",17,WHITE,True))],space_after=0)
text(s,Inches(L.inches+0.5),Inches(4.3),Inches(8.5),Inches(0.4),[P(R("無理な営業はいたしません。小さなお困りごとからお気軽にどうぞ。",13,RGBColor(0xF7,0xE4,0xD4)))],space_after=0)
text(s,Inches(L.inches+0.5),Inches(4.8),Inches(8.5),Inches(0.7),[P(R("📞 0120 - ○○○ - ○○○",30,WHITE,True))],space_after=0)
text(s,Inches(L.inches+0.5),Inches(5.45),Inches(8.5),Inches(0.4),[P(R("Web予約：www.tousaikan.com／reserve",13,RGBColor(0xF7,0xE4,0xD4)))],space_after=0)
x=L.inches
for txt in ["創業50年","施工55,000件超","自社施工・安心保証"]:
    w=chip(s,Inches(x),Inches(6.1),txt); x+=w.inches+0.2
text(s,L,Inches(6.85),Inches(12),Inches(0.4),
     [P(R("陶彩館（オーヤシマ株式会社）／※連絡先の○○部分は配布前に最新の公式情報へ差し替えてください。",9.5,INKSOFT))],space_after=0)

# ===================== 15. 種明かし（AIデモ）=====================
s = slide(INK)
rect(s,0,0,EW,EH,fill=RGBColor(0x2E,0x27,0x22))
rect(s,0,0,Inches(0.22),EH,fill=TERRA)
text(s,Inches(0.9),Inches(0.8),Inches(11),Inches(0.4),[P(R("AI TRAINING DEMO ｜ 種明かし",13,AMBER,True))],space_after=0)
text(s,Inches(0.9),Inches(1.5),Inches(11.5),Inches(1.6),
     [P(R("この14ページ、",30,WHITE,True),R("AI",30,AMBER,True),R("が御社サイトを学習して作りました。",30,WHITE,True))],space_after=2,line_spacing=1.25)
text(s,Inches(0.9),Inches(3.0),Inches(11.5),Inches(1.0),
     [P(R("会社情報・商材・強みを読み込ませるだけで、構成・デザイン・営業トークまで自動生成。研修を受ければ、御社の営業担当も“自分の言葉”でこのレベルの資料を作れるようになります。",15,RGBColor(0xD9,0xCF,0xC4)))],
     space_after=0,line_spacing=1.5)
demo=[("⏱️","数分","手作業なら数時間〜数日の資料作成を短縮"),
      ("🔁","再現可能","研修後は社員自身が何度でも作成・改善"),
      ("🎯","御社専用","公開情報・社内資料を学習し貴社仕様に")]
cw=Inches(3.75)
for i,(ic,t,b) in enumerate(demo):
    lx=Inches(0.9+i*4.0)
    rect(s,lx,Inches(4.55),cw,Inches(1.85),fill=RGBColor(0x3D,0x34,0x2D),shape=MSO_SHAPE.ROUNDED_RECTANGLE,round_=0.08)
    text(s,Inches(lx.inches+0.3),Inches(4.75),cw,Inches(0.5),[P(R(ic+"  ",20,WHITE),R(t,19,AMBER,True))],space_after=0)
    text(s,Inches(lx.inches+0.3),Inches(5.4),Inches(cw.inches-0.5),Inches(0.9),[P(R(b,12.5,RGBColor(0xD9,0xCF,0xC4)))],space_after=0,line_spacing=1.25)
text(s,Inches(0.9),Inches(6.75),Inches(11.5),Inches(0.4),
     [P(R("AI研修で、資料作成を“一部の人の特技”から“全員のスキル”へ。",13,AMBER,True))],space_after=0)

prs.save("/home/user/website-redesign-machine-/sales-materials/tousaikan/陶彩館_営業資料_AIデモ.pptx")
print("slides:", len(prs.slides._sldIdLst))
print("saved.")
