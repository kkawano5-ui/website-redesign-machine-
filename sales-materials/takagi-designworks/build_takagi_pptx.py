# -*- coding: utf-8 -*-
"""高城株式会社／株式会社デザインワークス様 AI活用構想 ご提案(叩き台) PPTX"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SLATE=RGBColor(0x28,0x38,0x45); BLUE=RGBColor(0x2f,0x6f,0x8f); GOLD=RGBColor(0xB9,0x85,0x2E)
PAPER=RGBColor(0xF6,0xF4,0xEF); CARD=RGBColor(0xFF,0xFF,0xFF); INK=RGBColor(0x2a,0x2f,0x34)
SOFT=RGBColor(0x6a,0x71,0x78); LINE=RGBColor(0xE3,0xDD,0xD2); CHIP=RGBColor(0xEC,0xE6,0xDA)
GREEN=RGBColor(0x4f,0x7a,0x5f); RED=RGBColor(0xC0,0x47,0x42); WHITE=RGBColor(0xFF,0xFF,0xFF)
SKY=RGBColor(0xE6,0xEF,0xF3)
JP="Meiryo"
EW,EH=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=EW; prs.slide_height=EH
BLANK=prs.slide_layouts[6]

def slide(bg=PAPER):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EW,EH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s
def _round(sp,v=0.08):
    try: sp.adjustments[0]=v
    except Exception: pass
def rect(s,l,t,w,h,fill=None,line=None,lw=1.0,shp=MSO_SHAPE.RECTANGLE,rnd=0.0):
    sp=s.shapes.add_shape(shp,l,t,w,h)
    if rnd:_round(sp,rnd)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def text(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=4,ls=1.0):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(2); tf.margin_top=tf.margin_bottom=Pt(2)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.line_spacing=ls
        for (txt,sz,col,bold) in para:
            r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name=JP; rPr=r.font._rPr
            for tag in ("a:latin","a:ea","a:cs"):
                el=rPr.find(qn(tag))
                if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
                el.set("typeface",JP)
    return tb
def P(*r): return list(r)
def R(t,s,c=INK,b=False): return (t,s,c,b)
def header(s,eyebrow,heading,page):
    rect(s,Inches(0.7),Inches(0.6),Inches(0.16),Inches(0.5),fill=GOLD)
    text(s,Inches(0.98),Inches(0.53),Inches(11),Inches(0.35),[P(R(eyebrow,12,GOLD,True))])
    text(s,Inches(0.95),Inches(0.84),Inches(11.7),Inches(0.7),[P(R(heading,25,SLATE,True))])
    rect(s,Inches(0.98),Inches(1.46),Inches(2.0),Inches(0.035),fill=BLUE)
    text(s,Inches(12.4),Inches(7.04),Inches(0.7),Inches(0.3),[P(R(str(page),11,SOFT))],align=PP_ALIGN.RIGHT)
    text(s,Inches(0.7),Inches(7.04),Inches(8),Inches(0.3),[P(R("高城株式会社／株式会社デザインワークス様 ご提案（叩き台）",9,SOFT))])
def card(s,l,t,w,h,title,body,icon=None,tcol=SLATE,fill=CARD,tsize=15,accent=None):
    rect(s,l,t,w,h,fill=fill,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.07)
    if accent: rect(s,l,t,w,Inches(0.12),fill=accent,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.4)
    yy=t.inches+0.2
    if icon:
        text(s,Inches(l.inches+0.25),Inches(yy),Inches(1),Inches(0.5),[P(R(icon,24,INK))],sa=0); yy+=0.5
    text(s,Inches(l.inches+0.25),Inches(yy),Inches(w.inches-0.45),Inches(0.6),[P(R(title,tsize,tcol,True))],sa=0,ls=1.05)
    text(s,Inches(l.inches+0.25),Inches(yy+(0.42 if icon else 0.5)),Inches(w.inches-0.45),h,[P(R(body,11.5,SOFT))],sa=0,ls=1.18)
L=Inches(0.7)

# ===== 1 COVER =====
s=slide(SLATE)
rect(s,0,0,EW,EH,fill=SLATE)
# blueprint-ish grid accent
for gx in range(0,14):
    rect(s,Inches(gx),0,Inches(0.006),EH,fill=RGBColor(0x32,0x44,0x52))
for gy in range(0,8):
    rect(s,0,Inches(gy),EW,Inches(0.006),fill=RGBColor(0x32,0x44,0x52))
rect(s,Inches(0.7),Inches(0.85),Inches(2.6),Inches(0.5),fill=None,line=GOLD,lw=1.5,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.3)
text(s,Inches(0.7),Inches(0.85),Inches(2.6),Inches(0.5),[P(R("AI活用 構想ご提案",12,GOLD,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
text(s,L,Inches(2.3),Inches(12.2),Inches(1.8),
     [P(R("眠っている“設計と顧客の情報”を、",34,WHITE,True)),P(R("AIで使える資産に変える。",34,WHITE,True))],sa=2,ls=1.2)
text(s,L,Inches(4.15),Inches(12),Inches(0.5),[P(R("図面・書類の電子化 ＋ 個人情報のAI移行・活用 ／ ご提案の叩き台",16,RGBColor(0xC9,0xD3,0xDA)))],sa=0)
text(s,L,Inches(5.25),Inches(12),Inches(0.5),[P(R("高城株式会社 ／ 株式会社デザインワークス　社長・部長　小早川 和廉 様",16,WHITE,True))],sa=0)
text(s,L,Inches(6.55),Inches(12),Inches(0.4),[P(R("2026年 ／ 方向性を一緒に固めるためのディスカッション資料",11.5,RGBColor(0x9A,0xA8,0xB2)))],sa=0)

# ===== 2 本日のゴール =====
s=slide()
header(s,"GOAL ｜ 本日のゴール","“正解を持ってくる”のではなく、一緒に決める",2)
text(s,L,Inches(1.95),Inches(11.8),Inches(0.7),
     [P(R("本資料は完成提案ではなく",13.5,SOFT),R("叩き台",13.5,GOLD,True),R("です。私たちの仮説をぶつけて、御社の実態に合う方向を本日その場で固めるのがゴールです。",13.5,SOFT))],sa=0,ls=1.4)
items=[("01","御社グループの理解をすり合わせる","建設・不動産＋設計、110-120名、すでにAI活用中という現在地"),
       ("02","社長の構想（④）を翻訳して確認","「設計と個人情報をAIで移行」を具体像に落とし、認識を合わせる"),
       ("03","進め方と“肝”を共有","データ移行の全体像と、個人情報・セキュリティの設計を提示"),
       ("04","本日決めたい論点を一緒に詰める","対象範囲・優先順位・PoCの第一歩をその場で合意したい")]
for i,(n,t,b) in enumerate(items):
    y=Inches(2.75+i*1.0)
    rect(s,L,y,Inches(11.9),Inches(0.88),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
    rect(s,Inches(0.9),Inches(y.inches+0.18),Inches(0.55),Inches(0.55),fill=SLATE,shp=MSO_SHAPE.OVAL)
    text(s,Inches(0.9),Inches(y.inches+0.18),Inches(0.55),Inches(0.55),[P(R(n,13,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(1.7),Inches(y.inches+0.12),Inches(10),Inches(0.4),[P(R(t,15.5,SLATE,True))],sa=0)
    text(s,Inches(1.7),Inches(y.inches+0.5),Inches(10),Inches(0.35),[P(R(b,11.5,SOFT))],sa=0)

# ===== 3 御社の理解 =====
s=slide()
header(s,"UNDERSTANDING ｜ 御社グループの理解","私たちが理解している、御社の現在地",3)
text(s,L,Inches(1.95),Inches(7.1),Inches(2.6),
     [P(R("高城株式会社（不動産）と株式会社デザインワークス（設計）が連携し、建設・不動産を一気通貫で手がけるグループ、と理解しています。",13,SOFT)),
      P(R("従業員110〜120名規模で、すでに一部社員がAIを活用し、公開APIも取り入れている——",13,SOFT),R("“使える組織”の土台はすでにある",13,SLATE,True),R("状態です。",13,SOFT)),
      P(R("次の一手は、現場の小技ではなく“設計×顧客情報”という会社の資産そのものをAIで活かすこと。",13.5,SLATE,True))],sa=8,ls=1.45)
for i,(v,u,lb) in enumerate([("110-120","名","グループ従業員"),("2","社","不動産 × 設計"),("活用中","","公開APIも導入済み"),("設計×個人情報","","次の主戦場（④）")]):
    col=i%2; row=i//2
    lx=Inches(8.2+col*2.45); ty=Inches(2.0+row*1.5)
    rect(s,lx,ty,Inches(2.25),Inches(1.35),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.1)
    sz=22 if len(v)>5 else 28
    text(s,lx,Inches(ty.inches+0.18),Inches(2.25),Inches(0.6),[P(R(v,sz,BLUE,True),R(u,13,BLUE,True))],align=PP_ALIGN.CENTER,sa=0)
    text(s,lx,Inches(ty.inches+0.85),Inches(2.25),Inches(0.4),[P(R(lb,11,SOFT))],align=PP_ALIGN.CENTER,sa=0)
text(s,L,Inches(5.2),Inches(11.9),Inches(0.3),[P(R("※ 公開情報からの推定を含みます。事実と異なる点は本日その場で修正させてください。",10,SOFT))],sa=0)

# ===== 4 ④の翻訳 =====
s=slide()
header(s,"VISION ｜ 社長の構想を翻訳する","「設計と個人情報をAIで移行」を、具体像に",4)
rect(s,L,Inches(1.95),Inches(11.9),Inches(1.15),fill=SLATE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
text(s,Inches(L.inches+0.4),Inches(2.12),Inches(11.1),Inches(0.9),
     [P(R("社長の言葉：",13,GOLD,True),R("「最終的には設計と個人情報の関係も、AIを使って移行したい」",14,WHITE,True))],sa=0,anchor=MSO_ANCHOR.MIDDLE,ls=1.2)
text(s,L,Inches(3.35),Inches(11.9),Inches(0.5),[P(R("私たちはこう受け止めました ▼（認識が違えば修正します）",13,SOFT,True))],sa=0)
cols=[("📂","いま","紙・PDFの図面、契約書、重要事項説明書、顧客台帳に、設計情報と個人情報が“バラバラ・属人的”に眠っている",CHIP,SLATE),
      ("🤖","AIで移行","AI-OCR＋情報抽出で、図面のメタ情報と顧客の個人情報を読み取り、構造化データに変換",SKY,BLUE),
      ("🗄️","これから","物件・設計・顧客が紐づいた検索可能なDBへ。「この土地の過去設計は？」に即答できる資産に",RGBColor(0xE7,0xF0,0xE9),GREEN)]
for i,(ic,t,b,bg,c) in enumerate(cols):
    lx=Inches(0.7+i*4.05)
    rect(s,lx,Inches(3.95),Inches(3.8),Inches(2.6),fill=bg,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
    text(s,Inches(lx.inches+0.25),Inches(4.12),Inches(3.3),Inches(0.5),[P(R(ic+"  "+t,15,c,True))],sa=0)
    text(s,Inches(lx.inches+0.25),Inches(4.75),Inches(3.35),Inches(1.7),[P(R(b,11.5,INK))],sa=0,ls=1.3)
    if i<2: text(s,Inches(lx.inches+3.72),Inches(4.9),Inches(0.5),Inches(0.5),[P(R("▶",18,GOLD,True))],align=PP_ALIGN.CENTER,sa=0)

# ===== 5 なぜ今か =====
s=slide()
header(s,"WHY NOW ｜ なぜ今やる価値があるか","“眠った資産”は、年々取り出しにくくなる",5)
text(s,L,Inches(1.95),Inches(12),Inches(0.5),[P(R("建設・不動産の情報は、年月とともに紙・担当者の頭の中に固定化します。デジタル資産化を先送りするほどコストが上がります。",13,SOFT))],sa=0)
pains=[("🔍","探せない・即答できない","過去物件・図面・顧客履歴が探せず、問い合わせや再提案に時間がかかる"),
       ("👤","属人化・退職リスク","“あの人しか知らない”状態。ベテラン退職で情報ごと失われる"),
       ("📑","紙・劣化リスク","図面・契約書の紙原本は劣化・紛失・災害に弱い"),
       ("⚖️","個人情報の管理不安","どこに誰の情報があるか把握しきれず、漏えい・コンプラのリスク"),
       ("🔁","再利用できない","せっかくの設計ノウハウが次の提案・営業に活かしきれていない"),
       ("⏱️","対応スピード","相続・売買・リフォーム相談での即応力が、そのまま受注力になる")]
cw=Inches(3.78); ch=Inches(1.55)
for i,(ic,t,b) in enumerate(pains):
    col=i%3; row=i//3
    lx=Inches(0.7+col*4.05); ty=Inches(2.6+row*1.8)
    rect(s,lx,ty,cw,ch,fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.07)
    text(s,Inches(lx.inches+0.25),Inches(ty.inches+0.2),Inches(1),Inches(0.4),[P(R(ic,20,INK))],sa=0)
    text(s,Inches(lx.inches+0.95),Inches(ty.inches+0.24),Inches(cw.inches-1.1),Inches(0.5),[P(R(t,14,SLATE,True))],sa=0,ls=1.05)
    text(s,Inches(lx.inches+0.25),Inches(ty.inches+0.82),Inches(cw.inches-0.45),Inches(0.7),[P(R(b,11,SOFT))],sa=0,ls=1.15)

# ===== 6 全体像 パイプライン =====
s=slide()
header(s,"SOLUTION ｜ 構想の全体像","図面・書類 → AI抽出 → 使えるDB",6)
text(s,L,Inches(1.9),Inches(12),Inches(0.5),[P(R("レガシー資料を“読み取り→構造化→紐付け”する4ステップ。AIが人手の入力・転記を肩代わりします。",12.5,SOFT))],sa=0)
steps=[("①","取り込む","紙図面・PDF・契約書・顧客台帳をスキャン／集約",CHIP,SLATE),
       ("②","AIで読む","AI-OCR＋抽出で、図面メタ情報と個人情報をテキスト化・項目化",SKY,BLUE),
       ("③","個人情報を守る","氏名・住所等を自動検出し、マスキング／アクセス権限を付与",RGBColor(0xF3,0xE7,0xE6),RED),
       ("④","DBへ紐付け","物件 × 設計 × 顧客 を関連づけ、検索可能なデータベースへ",RGBColor(0xE7,0xF0,0xE9),GREEN)]
bw=Inches(2.78)
for i,(n,t,b,bg,c) in enumerate(steps):
    lx=Inches(0.62+i*2.95); ty=Inches(2.55)
    rect(s,lx,ty,bw,Inches(2.5),fill=bg,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
    rect(s,lx,ty,bw,Inches(0.65),fill=c,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.12)
    text(s,lx,ty,bw,Inches(0.65),[P(R(n+"  "+t,14,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(lx.inches+0.22),Inches(ty.inches+0.85),Inches(bw.inches-0.4),Inches(1.5),[P(R(b,11.5,INK))],sa=0,ls=1.3)
    if i<3: text(s,Inches(lx.inches+bw.inches-0.02),Inches(ty.inches+0.95),Inches(0.4),Inches(0.5),[P(R("▶",15,GOLD,True))],align=PP_ALIGN.CENTER,sa=0)
rect(s,L,Inches(5.4),Inches(11.9),Inches(1.05),fill=CHIP,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
text(s,Inches(L.inches+0.3),Inches(5.56),Inches(11.3),Inches(0.8),
     [P(R("出口イメージ：",12.5,SLATE,True),R("「○○町の過去設計を見せて」「この顧客の取引履歴は？」に、社員の誰もが数秒で答えられる状態。",12.5,INK))],sa=0,ls=1.3)

# ===== 7 個人情報・セキュリティ（肝） =====
s=slide()
header(s,"SECURITY ｜ ここが肝","個人情報を“守りながら”使う設計",7)
text(s,L,Inches(1.92),Inches(12),Inches(0.6),
     [P(R("建設・不動産は契約書・図面に個人情報が詰まっています。AI活用の可否は、ここの設計で決まります。先に固めます。",13,SOFT))],sa=0)
secs=[("🔒","学習されないAIを選定","入力データをモデル学習に使わない商用APIを採用。情報が外部学習に流れない契約形態で運用"),
      ("🎭","送信前にマスキング","氏名・住所・連絡先など個人情報を自動検出し、必要に応じ伏せてから処理。最小限の情報だけ扱う"),
      ("🗝️","アクセス権限・監査ログ","誰がどの情報を見たか記録。役職・部門単位で閲覧範囲を制御し、内部統制を担保"),
      ("🏠","保管場所の選択","社内・閉域・国内クラウドなど、御社のポリシーに合わせて保管先と経路を設計"),
      ("📜","法令順守","個人情報保護法・宅建業法の要件に沿った取得・利用・保管・廃棄のルール整備"),
      ("🧪","小さく検証","まず限定データでPoC。安全性を確認しながら段階的に範囲を広げる")]
cw=Inches(3.78); ch=Inches(1.62)
for i,(ic,t,b) in enumerate(secs):
    col=i%3; row=i//3
    lx=Inches(0.7+col*4.05); ty=Inches(2.62+row*1.82)
    rect(s,lx,ty,cw,ch,fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.07)
    rect(s,lx,ty,Inches(0.12),ch,fill=RED,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.4)
    text(s,Inches(lx.inches+0.28),Inches(ty.inches+0.18),Inches(1),Inches(0.4),[P(R(ic,18,INK))],sa=0)
    text(s,Inches(lx.inches+0.95),Inches(ty.inches+0.22),Inches(cw.inches-1.1),Inches(0.5),[P(R(t,13.5,SLATE,True))],sa=0,ls=1.0)
    text(s,Inches(lx.inches+0.28),Inches(ty.inches+0.8),Inches(cw.inches-0.5),Inches(0.75),[P(R(b,10.8,SOFT))],sa=0,ls=1.15)

# ===== 8 進め方 =====
s=slide()
header(s,"ROADMAP ｜ 進め方","小さく試して、確実に広げる",8)
steps=[("PHASE 1","PoC（実証）","対象を1種類の書類・一部物件に絞り、抽出精度と安全性を検証。効果と勘所を掴む","〜1〜2ヶ月",BLUE),
       ("PHASE 2","本格移行","精度の出た書類から順に電子化・DB化。個人情報の運用ルールを確立し定着","〜3〜6ヶ月",SLATE),
       ("PHASE 3","活用・横展開","検索・再提案・営業活用へ。デザインワークスの設計ノウハウ活用、グループ全体へ","継続",GOLD)]
for i,(ph,t,b,d,c) in enumerate(steps):
    lx=Inches(0.7+i*4.05)
    rect(s,lx,Inches(2.2),Inches(3.8),Inches(3.5),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
    rect(s,lx,Inches(2.2),Inches(3.8),Inches(0.95),fill=c,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.12)
    text(s,lx,Inches(2.33),Inches(3.8),Inches(0.4),[P(R(ph,15,WHITE,True))],align=PP_ALIGN.CENTER,sa=0)
    text(s,lx,Inches(2.73),Inches(3.8),Inches(0.35),[P(R(d,11.5,RGBColor(0xEC,0xEC,0xEC),True))],align=PP_ALIGN.CENTER,sa=0)
    text(s,Inches(lx.inches+0.3),Inches(3.4),Inches(3.2),Inches(0.5),[P(R(t,16,SLATE,True))],sa=0)
    text(s,Inches(lx.inches+0.3),Inches(4.05),Inches(3.2),Inches(1.5),[P(R(b,12,SOFT))],sa=0,ls=1.3)
    if i<2: text(s,Inches(lx.inches+3.72),Inches(3.6),Inches(0.4),Inches(0.5),[P(R("▶",18,GOLD,True))],align=PP_ALIGN.CENTER,sa=0)
text(s,L,Inches(6.0),Inches(12),Inches(0.5),
     [P(R("いきなり全部はやりません。PHASE 1で“効くと分かってから”投資判断ができる進め方です。",13,SLATE,True))],sa=0)

# ===== 9 広がり =====
s=slide()
header(s,"EXPANSION ｜ 広がり","同じ仕組みが、他業務にも効く",9)
more=[("📝","物件紹介文の自動生成","物件データから広告・ポータル掲載文を自動作成"),
      ("📋","重要事項説明書・契約書ドラフト","ひな型＋物件情報から下書きを自動生成、確認は人が行う"),
      ("📐","図面からの情報抽出","図面から面積・部屋数・仕様を読み取り、台帳へ自動入力"),
      ("💬","顧客問い合わせ一次対応","過去物件・FAQをもとに、問い合わせへ即時ドラフト回答"),
      ("🏗️","設計ナレッジの再利用","過去設計を検索し、新規提案・見積のたたき台に活用"),
      ("📊","社内ナレッジ横断検索","設計×顧客×物件を横断して“聞けば返ってくる”状態に")]
cw=Inches(3.78); ch=Inches(1.62)
for i,(ic,t,b) in enumerate(more):
    col=i%3; row=i//3
    lx=Inches(0.7+col*4.05); ty=Inches(2.1+row*1.85)
    card(s,lx,ty,cw,ch,t,b,icon=ic,tsize=13.5)
text(s,L,Inches(5.95),Inches(12),Inches(0.4),
     [P(R("④のデータ基盤ができると、これらは“その上のアプリ”として次々に載せられます。",12.5,SLATE,True))],sa=0)

# ===== 10 本日決めたいこと（叩き台の核）=====
s=slide(SLATE)
rect(s,0,0,EW,EH,fill=SLATE)
text(s,L,Inches(0.7),Inches(11),Inches(0.4),[P(R("DISCUSSION ｜ 本日、一緒に決めたいこと",13,GOLD,True))],sa=0)
text(s,L,Inches(1.25),Inches(12),Inches(0.7),[P(R("この5つが決まれば、次回は“動くもの”を持ってこられます。",24,WHITE,True))],sa=0)
qs=[("対象","まず何の書類・どの物件群から始めるか？（一番困っている / 一番数が多い）"),
    ("現状","図面・顧客情報は今どんな形式か？（紙・PDF・既存システム・Excel）"),
    ("個人情報","社内の取扱いルール・NGライン・保管ポリシーは？クラウド利用の可否は？"),
    ("ゴール","“移行”の先に何をしたいか？（検索・再提案・営業・相続対応…優先順位）"),
    ("体制","推進担当は誰か？デザインワークスと高城のどちらから着手するか？")]
for i,(k,q) in enumerate(qs):
    y=Inches(2.35+i*0.92)
    rect(s,L,y,Inches(11.9),Inches(0.8),fill=RGBColor(0x31,0x43,0x51),shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.1)
    rect(s,Inches(0.9),Inches(y.inches+0.14),Inches(1.7),Inches(0.52),fill=GOLD,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.3)
    text(s,Inches(0.9),Inches(y.inches+0.14),Inches(1.7),Inches(0.52),[P(R(k,12.5,SLATE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    text(s,Inches(2.8),y,Inches(9.0),Inches(0.8),[P(R(q,12.5,WHITE))],anchor=MSO_ANCHOR.MIDDLE,sa=0,ls=1.2)

# ===== 11 体制・効果 =====
s=slide()
header(s,"VALUE ｜ 体制と効果イメージ","“作って終わり”にしない伴走",11)
for i,(v,lb,c) in enumerate([("属人化 → 資産化","個人の頭の中の情報を、会社の検索可能な資産へ",BLUE),
                             ("探す時間 → ゼロ近く","過去物件・図面・顧客履歴に即アクセス",GREEN),
                             ("守りながら活用","個人情報を保護しつつAIで使う体制を確立",RED),
                             ("内製でコスト最適","研修で社内に運用力を残し、外注依存を避ける",GOLD)]):
    col=i%2; row=i//2
    lx=Inches(0.7+col*6.05); ty=Inches(2.1+row*1.5)
    rect(s,lx,ty,Inches(5.85),Inches(1.3),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.07)
    rect(s,lx,ty,Inches(0.14),Inches(1.3),fill=c,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.4)
    text(s,Inches(lx.inches+0.35),Inches(ty.inches+0.2),Inches(5.3),Inches(0.5),[P(R(v,16,SLATE,True))],sa=0)
    text(s,Inches(lx.inches+0.35),Inches(ty.inches+0.72),Inches(5.3),Inches(0.5),[P(R(lb,12,SOFT))],sa=0,ls=1.15)
rect(s,L,Inches(5.25),Inches(11.9),Inches(1.2),fill=CHIP,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
text(s,Inches(L.inches+0.3),Inches(5.45),Inches(11.3),Inches(0.9),
     [P(R("私たちの関わり方：",13,SLATE,True),R("構想設計／PoC伴走／個人情報・セキュリティ設計／社員向けAI研修。御社の“使える組織”の土台に、武器を載せます。",12.5,INK))],sa=0,ls=1.35)

# ===== 12 次のアクション =====
s=slide()
header(s,"NEXT ｜ 次のアクション","まずは“1種類の書類”から、一緒に。",12)
for i,(t,b) in enumerate([("① 対象を1つ決める","本日の議論で、PoC対象の書類・物件群を1つ選定"),
                          ("② サンプルで実証（PoC）","実データ一部で抽出精度と安全性を検証。効果を数字で確認"),
                          ("③ 移行＋研修を並走","DB化を進めつつ、社員が自分で使える研修を実施")]):
    y=Inches(2.1+i*1.0)
    rect(s,L,y,Inches(8.5),Inches(0.85),fill=CARD,line=LINE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.1)
    text(s,Inches(L.inches+0.3),Inches(y.inches+0.12),Inches(8),Inches(0.4),[P(R(t,15,SLATE,True))],sa=0)
    text(s,Inches(L.inches+0.3),Inches(y.inches+0.5),Inches(8),Inches(0.35),[P(R(b,11.5,SOFT))],sa=0)
rect(s,Inches(9.5),Inches(2.1),Inches(3.1),Inches(2.75),fill=SLATE,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.06)
text(s,Inches(9.5),Inches(2.5),Inches(3.1),Inches(0.5),[P(R("ご一緒できること",13,GOLD,True))],align=PP_ALIGN.CENTER,sa=0)
text(s,Inches(9.72),Inches(3.1),Inches(2.7),Inches(1.6),
     [P(R("✓ 構想・要件の整理",12,WHITE,True)),P(R("✓ PoC設計・実証",12,WHITE,True)),P(R("✓ 個人情報・安全設計",12,WHITE,True)),P(R("✓ 社員向けAI研修",12,WHITE,True))],sa=7,ls=1.15)
rect(s,L,Inches(5.4),Inches(11.9),Inches(1.05),fill=CHIP,shp=MSO_SHAPE.ROUNDED_RECTANGLE,rnd=0.05)
text(s,Inches(L.inches+0.3),Inches(5.6),Inches(11.3),Inches(0.7),
     [P(R("“設計と個人情報をAIで移行したい”——その第一歩を、リスクの小さいPoCから一緒に踏み出しましょう。",13.5,SLATE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)

prs.save("/home/user/website-redesign-machine-/sales-materials/takagi-designworks/高城_デザインワークス様_AI構想提案_叩き台.pptx")
print("slides:",len(prs.slides._sldIdLst))
