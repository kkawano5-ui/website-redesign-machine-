# -*- coding: utf-8 -*-
"""提案資料 共通デザイン部品（中小企業版・大企業版で共有）。
THEME定数とヘルパー、render() を提供する。各版は build_*.py でスライドを定義して render() を呼ぶ。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- THEME ----------
NAVY   = RGBColor(0x14,0x22,0x3D)
NAVY2  = RGBColor(0x1C,0x32,0x59)
BLUE   = RGBColor(0x2A,0x5A,0xCC)
BLUE_D = RGBColor(0x1D,0x3B,0x6E)
BLUESOFT=RGBColor(0xEC,0xF1,0xFA)
GOLD   = RGBColor(0xB4,0x8A,0x3A)
INK    = RGBColor(0x1B,0x22,0x30)
SUB    = RGBColor(0x57,0x5F,0x70)
FAINT  = RGBColor(0x76,0x7C,0x87)
LINE   = RGBColor(0xE3,0xE6,0xEC)
LINE2  = RGBColor(0xED,0xEF,0xF3)
PAPER  = RGBColor(0xFF,0xFF,0xFF)
PANEL  = RGBColor(0xF5,0xF7,0xFA)
DELTA  = RGBColor(0xB0,0x4E,0x2E)
DELTASOFT=RGBColor(0xF6,0xEC,0xE6)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
CREAM  = RGBColor(0xD8,0xE2,0xF2)

EA   = "游ゴシック"
EAL  = "Yu Gothic"
EW = 12192000
EH = 6858000
MX = Inches(0.92)
CW = EW - 2*int(MX)

def new_prs():
    prs = Presentation(); prs.slide_width = EW; prs.slide_height = EH
    return prs

def slide(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, EW, EH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit=False
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _set_font(run, size, color, bold=False, ea=EA, latin=EAL, italic=False, spacing=None):
    f = run.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea','a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', ea)
    if spacing is not None: rPr.set('spc', str(int(spacing*100)))

def textbox(s, x,y,w,h, spans, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, line=1.3, ea=EA, space_after=0, wrap=True):
    tb = s.shapes.add_textbox(x,y,w,h); tf = tb.text_frame; tf.word_wrap = wrap
    try: tf.auto_size = None
    except Exception: pass
    tf.vertical_anchor = anchor
    for m in ('left','right','top','bottom'): setattr(tf, 'margin_'+m, 0)
    if isinstance(spans, str): spans = [spans]
    first = True
    for para in spans:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align
        try: p.line_spacing = line
        except Exception: pass
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        runs = para if isinstance(para, list) else [(para, {})]
        for txt, kw in ([(para,{})] if isinstance(para,str) else runs):
            run = p.add_run(); run.text = txt
            _set_font(run, kw.get('size',size), kw.get('color',color), kw.get('bold',bold),
                      kw.get('ea',ea), kw.get('latin',EAL), kw.get('italic',False), kw.get('spacing'))
    return tb

def rect(s, x,y,w,h, fill=None, line_color=None, line_w=1.0, rounded=False, radius=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x,y,w,h)
    if rounded:
        try: shp.adjustments[0]=radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None: shp.line.fill.background()
    else: shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def hline(s, x,y,w, color=LINE, weight=1.0):
    ln = s.shapes.add_connector(2, x,y, x+w,y)
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit=False
    return ln

def kicker(s, txt, x=MX, y=Inches(0.62), color=BLUE, dark=False):
    c = CREAM if dark else color
    rect(s, x, y+Pt(6), Inches(0.30), Pt(2.2), fill=c)
    textbox(s, x+Inches(0.42), y, Inches(9), Inches(0.3),
            [[(txt,{'size':12.5,'color':c,'bold':True,'spacing':1.6})]], anchor=MSO_ANCHOR.MIDDLE)

def title(s, spans, y=Inches(1.02), size=30, color=INK, w=None, line=1.28):
    textbox(s, MX, y, w or CW, Inches(1.6),
            spans if isinstance(spans,list) else [[(spans,{'size':size,'color':color,'bold':True})]],
            size=size, color=color, bold=True, line=line)

def footer(s, idx, total, sec, dark=False):
    c = WHITE if dark else INK
    cf = RGBColor(0xAF,0xBC,0xD4) if dark else FAINT
    y = EH - Inches(0.46)
    rect(s, MX, y+Pt(4), Inches(0.11), Inches(0.11), fill=BLUE if not dark else CREAM)
    textbox(s, MX+Inches(0.20), y, Inches(4), Inches(0.3),
            [[('株式会社AIスキル',{'size':10.5,'color':c,'bold':True})]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, EW-MX-Inches(4), y, Inches(4), Inches(0.3),
            [[('%02d / %02d   %s'%(idx,total,sec),{'size':10,'color':cf,'latin':'Arial'})]],
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def logo(s, x=MX, y=Inches(0.6), dark=True):
    c = WHITE if dark else INK
    d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y+Pt(2), Inches(0.20), Inches(0.20))
    d.fill.solid(); d.fill.fore_color.rgb = CREAM if dark else BLUE; d.line.fill.background(); d.shadow.inherit=False
    textbox(s, x+Inches(0.32), y-Pt(2), Inches(4), Inches(0.4),
            [[('AIスキル',{'size':17,'color':c,'bold':True})]], anchor=MSO_ANCHOR.MIDDLE)

def pill_row(s, items, x, y, fill=BLUESOFT, line_color=RGBColor(0xDD,0xE6,0xF6), tc=BLUE_D,
             size=12, h=0.42, gap=0.16, maxx=None, char_w=0.155, pad=0.30):
    """Render pills, wrapping at maxx. Returns final y."""
    maxx = maxx or (EW-MX); px=x; py=y
    for p in items:
        w=Inches(pad+len(p)*char_w)
        if px+w>maxx: px=x; py+=Inches(h+0.13)
        rect(s,px,py,w,Inches(h),fill=fill,line_color=line_color,rounded=True,radius=0.5)
        textbox(s,px,py,w,Inches(h),[[(p,{'size':size,'color':tc,'bold':True})]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        px+=w+Inches(gap)
    return py

def render(prs, slides, out):
    total=len(slides)
    for i,fn in enumerate(slides, start=1):
        fn(prs, i, total)
    prs.save(out)
    print("saved", out, "slides:", total)
    return out
