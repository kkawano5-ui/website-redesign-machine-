# -*- coding: utf-8 -*-
"""法人AI研修 提案資料（PPT版）ビルダー。デザイン定数はTHEMEで一元管理。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

# ---------- THEME ----------
NAVY   = RGBColor(0x14,0x22,0x3D)
NAVY2  = RGBColor(0x1C,0x32,0x59)
BLUE   = RGBColor(0x2A,0x5A,0xCC)
BLUE_D = RGBColor(0x1D,0x3B,0x6E)
BLUESOFT=RGBColor(0xEC,0xF1,0xFA)
GOLD   = RGBColor(0xB4,0x8A,0x3A)
INK    = RGBColor(0x1B,0x22,0x30)
SUB    = RGBColor(0x57,0x5F,0x70)
FAINT  = RGBColor(0x97,0x9E,0xAB)
LINE   = RGBColor(0xE3,0xE6,0xEC)
LINE2  = RGBColor(0xED,0xEF,0xF3)
PAPER  = RGBColor(0xFF,0xFF,0xFF)
PANEL  = RGBColor(0xF5,0xF7,0xFA)
DELTA  = RGBColor(0xB0,0x4E,0x2E)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
CREAM  = RGBColor(0xD8,0xE2,0xF2)

EA   = "游ゴシック"          # East-Asian font (Yu Gothic) — premium, ships with Office/Win/Mac
EAL  = "Yu Gothic"          # latin name of the same family
EW = 12192000  # 13.333in in EMU
EH = 6858000   # 7.5in
MX = Inches(0.92)  # left/right margin
CW = EW - 2*int(MX)  # content width

prs = Presentation()
prs.slide_width = EW
prs.slide_height = EH
BLANK = prs.slide_layouts[6]

def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, EW, EH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _set_font(run, size, color, bold=False, ea=EA, latin=EAL, italic=False, spacing=None):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea','a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', ea)
    if spacing is not None:
        rPr.set('spc', str(int(spacing*100)))

def textbox(s, x,y,w,h, spans, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, line=1.3, ea=EA, space_after=0, wrap=True):
    """spans: str OR list of paragraphs; each paragraph is str or list of (text, kwargs)."""
    tb = s.shapes.add_textbox(x,y,w,h); tf = tb.text_frame
    tf.word_wrap = wrap
    try: tf.auto_size = None
    except Exception: pass
    tf.vertical_anchor = anchor
    for m in ('left','right','top','bottom'):
        setattr(tf, 'margin_'+m, 0)
    if isinstance(spans, str): spans = [spans]
    first = True
    for para in spans:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        try: p.line_spacing = line
        except Exception: pass
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        runs = para if isinstance(para, list) else [(para, {})]
        for txt, kw in ([ (para,{}) ] if isinstance(para,str) else runs):
            run = p.add_run(); run.text = txt
            _set_font(run, kw.get('size',size), kw.get('color',color),
                      kw.get('bold',bold), kw.get('ea',ea), kw.get('latin',EAL),
                      kw.get('italic',False), kw.get('spacing'))
    return tb

def rect(s, x,y,w,h, fill=None, line_color=None, line_w=1.0, rounded=False, radius=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x,y,w,h)
    if rounded:
        try: shp.adjustments[0]=radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None: shp.line.fill.background()
    else: shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def hline(s, x,y,w, color=LINE, weight=1.0):
    ln = s.shapes.add_connector(2, x,y, x+w,y)
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit=False
    return ln

def kicker(s, txt, x=MX, y=Inches(0.62), color=BLUE, dark=False):
    c = CREAM if dark else color
    rect(s, x, y+Pt(6), Inches(0.30), Pt(2.2), fill=c)
    textbox(s, x+Inches(0.42), y, Inches(8), Inches(0.3),
            [[(txt,{'size':12.5,'color':c,'bold':True,'spacing':1.6})]], anchor=MSO_ANCHOR.MIDDLE)

def title(s, spans, y=Inches(1.02), size=30, color=INK, w=None, line=1.28):
    textbox(s, MX, y, w or CW, Inches(1.5),
            spans if isinstance(spans,list) else [[(spans,{'size':size,'color':color,'bold':True})]],
            size=size, color=color, bold=True, line=line)

def footer(s, idx, total, sec, dark=False):
    c = RGBColor(0xFF,0xFF,0xFF) if dark else INK
    cf = RGBColor(0xAF,0xBC,0xD4) if dark else FAINT
    y = EH - Inches(0.46)
    rect(s, MX, y+Pt(4), Inches(0.11), Inches(0.11), fill=BLUE if not dark else CREAM)
    textbox(s, MX+Inches(0.20), y, Inches(4), Inches(0.3),
            [[('株式会社AIスキル',{'size':10.5,'color':c,'bold':True})]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, EW-MX-Inches(4), y, Inches(4), Inches(0.3),
            [[('%02d / %02d   %s'%(idx,total,sec),{'size':10,'color':cf,'latin':'Arial'})]],
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def logo(s, x=MX, y=Inches(0.6), dark=True):
    c = WHITE if dark else INK
    d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y+Pt(2), Inches(0.20), Inches(0.20))
    d.fill.solid(); d.fill.fore_color.rgb = CREAM if dark else BLUE; d.line.fill.background(); d.shadow.inherit=False
    textbox(s, x+Inches(0.32), y-Pt(2), Inches(4), Inches(0.4),
            [[('AIスキル',{'size':17,'color':c,'bold':True})]], anchor=MSO_ANCHOR.MIDDLE)

SLIDES=[]
def register(fn): SLIDES.append(fn); return fn

# ============================================================ SLIDES
@register
def s_cover(s, idx, total):
    s = slide(NAVY)
    logo(s, MX, Inches(0.62), dark=True)
    textbox(s, MX, Inches(2.35), Inches(11.0), Inches(2.2),
        [[('研修で終わらせない。',{'size':44,'color':WHITE,'bold':True})],
         [('“現場が変わる”',{'size':44,'color':CREAM,'bold':True}),('まで。',{'size':44,'color':WHITE,'bold':True})]],
        line=1.3)
    textbox(s, MX, Inches(4.3), Inches(10.8), Inches(1.4),
        [[('「使い方を教える研修」ではありません。',{'size':15.5,'color':CREAM})],
         [('受講後に“あの業務が、実際に変わっている”状態まで設計してお渡しする、',{'size':15.5,'color':CREAM})],
         [('実践型の法人向け生成AI研修のご提案です。',{'size':15.5,'color':CREAM})]], line=1.6)
    # meta
    yb=Inches(5.95)
    for i,(k,v) in enumerate([('ご提案先','◯◯◯◯株式会社 御中'),('ご提案','株式会社AIスキル'),('ご提案日','2026年6月')]):
        xx = MX+Inches(i*3.45)
        textbox(s, xx, yb, Inches(3.2), Inches(0.3),[[(k,{'size':10.5,'color':RGBColor(0x8F,0xA3,0xC6)})]])
        textbox(s, xx, yb+Inches(0.28), Inches(3.3), Inches(0.4),[[(v,{'size':14,'color':WHITE,'bold':True})]])
    rect(s, MX, Inches(2.0), Inches(0.55), Pt(3), fill=GOLD)
    footer(s, idx, total, '表紙', dark=True)

@register
def s_company(s, idx, total):
    s = slide(PAPER)
    kicker(s,'会社紹介')
    title(s,[[('AIで、人を、社会を、',{'size':30,'color':INK,'bold':True}),('未来を変える。',{'size':30,'color':BLUE,'bold':True})]])
    textbox(s, MX, Inches(1.95), Inches(7.4), Inches(1.4),
        [[('最先端のAIスキルを、体系的・実践的に学べる教育環境を提供。',{'size':14.5,'color':SUB})],
         [('AI導入から運用・定着までを、一貫して支援しています。',{'size':14.5,'color':SUB})],
         [('“使いこなせる人材”を育てることを、最も大切にしています。',{'size':14.5,'color':SUB})]], line=1.7)
    # business pills
    textbox(s, MX, Inches(3.95), Inches(4), Inches(0.3),[[('提供サービス',{'size':11,'color':FAINT,'spacing':0.8})]])
    pills=['AI活用サポート事業','教育事業','コミュニティ事業','SNS事業']
    px=MX; py=Inches(4.32)
    for p in pills:
        w=Inches(0.30+len(p)*0.16)
        rect(s,px,py,w,Inches(0.42),fill=BLUESOFT,line_color=RGBColor(0xDD,0xE6,0xF6),rounded=True,radius=0.5)
        textbox(s,px,py,w,Inches(0.42),[[(p,{'size':12,'color':BLUE_D,'bold':True})]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        px+=w+Inches(0.18)
    # right info card
    cx=EW-MX-Inches(3.8)
    rect(s,cx,Inches(1.85),Inches(3.8),Inches(3.05),fill=PANEL,line_color=LINE,rounded=True,radius=0.05)
    info=[('会社名','株式会社AIスキル'),('代表者','立石 亮介'),('設立','2024年12月'),
          ('所在地','東京都渋谷区神南1-11-4'),('事業','AI活用支援／教育／コミュニティ')]
    iy=Inches(2.15)
    for k,v in info:
        textbox(s,cx+Inches(0.32),iy,Inches(1.1),Inches(0.4),[[(k,{'size':10.5,'color':FAINT})]])
        textbox(s,cx+Inches(1.35),iy,Inches(2.3),Inches(0.5),[[(v,{'size':12.5,'color':INK,'bold':True})]],line=1.2)
        iy+=Inches(0.56)
    footer(s, idx, total, '会社紹介')

@register
def s_track(s, idx, total):
    s = slide(PAPER)
    kicker(s,'実績')
    title(s,[[('“教える”だけでなく“成果を出す”研修を、',{'size':28,'color':INK,'bold':True})],
             [('数多くの企業・経営者に。',{'size':28,'color':INK,'bold':True})]],y=Inches(1.0),line=1.25)
    stats=[('200','社超','研修・セミナー\n導入企業'),
           ('150,000','名超','のべセミナー\n参加者数'),
           ('15,000','名超','AI学習コミュニティ\n会員数'),
           ('3,000','名超','経営者が\n研修・セミナー受講'),
           ('700','本超','自社制作の\nAI動画教材')]
    n=len(stats); gap=Inches(0.24)
    cw=(CW-gap*(n-1))/n
    y=Inches(2.55); h=Inches(2.3)
    for i,(num,unit,lab) in enumerate(stats):
        x=MX+i*(cw+gap)
        rect(s,x,y,cw,h,fill=WHITE,line_color=LINE,rounded=True,radius=0.06)
        nsize = 26 if len(num)>=6 else 37
        textbox(s,x,y+Inches(0.42),cw,Inches(0.95),
            [[(num,{'size':nsize,'color':BLUE_D,'bold':True,'latin':'Arial'}),(unit,{'size':12.5,'color':BLUE,'bold':True})]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        rect(s,x+cw/2-Inches(0.18),y+Inches(1.42),Inches(0.36),Pt(2),fill=GOLD)
        textbox(s,x,y+Inches(1.55),cw,Inches(0.7),
            [[(lab.replace('\n',''),{'size':11.5,'color':SUB})]] if '\n' not in lab else
            [[(lab.split('\n')[0],{'size':11.5,'color':SUB})],[(lab.split('\n')[1],{'size':11.5,'color':SUB})]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,line=1.35)
    textbox(s,MX,Inches(5.25),CW,Inches(0.5),
        [[('上場企業から数名規模の会社まで、全国・全業種で導入。',{'size':14,'color':INK,'bold':True}),
          ('  業種・規模を問わず成果を出しています。',{'size':14,'color':SUB})]])
    footer(s, idx, total, '実績')

@register
def s_supervisor(s, idx, total):
    s = slide(PAPER)
    kicker(s,'監修・体制')
    title(s,[[('研修屋ではなく、',{'size':29,'color':INK,'bold':True}),('AIの専門性と事業の両方',{'size':29,'color':BLUE,'bold':True}),('を持つチーム。',{'size':29,'color':INK,'bold':True})]])
    # photo placeholder
    px=MX; py=Inches(2.1)
    rect(s,px,py,Inches(2.7),Inches(2.7),fill=PANEL,line_color=LINE,rounded=True,radius=0.05)
    textbox(s,px,py+Inches(2.15),Inches(2.7),Inches(0.4),[[('監修者',{'size':12,'color':FAINT,'bold':True})]],align=PP_ALIGN.CENTER)
    bx=px+Inches(3.1)
    textbox(s,bx,Inches(2.1),Inches(7),Inches(0.4),[[('研修プログラム監修',{'size':11,'color':BLUE,'bold':True,'spacing':1.2})]])
    textbox(s,bx,Inches(2.45),Inches(7),Inches(0.6),[[('鈴木 章央',{'size':24,'color':INK,'bold':True}),('  Suzuki Akihiro',{'size':13,'color':FAINT})]])
    textbox(s,bx,Inches(3.2),Inches(7.3),Inches(1.6),
        [[('国立大学大学院でAI（人工知能）を研究し、修士・博士号を取得。2017年に',{'size':13.5,'color':SUB})],
         [('AI開発企業を創業し、受託開発・技術顧問・教育事業を展開。企業のAI導入支援を',{'size':13.5,'color':SUB})],
         [('多数手がける。AIプログラミング講座の講師歴は5年（Udemyでも講座を展開）。',{'size':13.5,'color':SUB})]],line=1.7)
    pills=['AI研究 修士・博士','AI受託開発・技術顧問','講師歴5年']
    qx=bx
    for p in pills:
        w=Inches(0.3+len(p)*0.15)
        rect(s,qx,Inches(4.55),w,Inches(0.4),fill=BLUESOFT,rounded=True,radius=0.5)
        textbox(s,qx,Inches(4.55),w,Inches(0.4),[[(p,{'size':11,'color':BLUE_D,'bold':True})]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        qx+=w+Inches(0.16)
    footer(s, idx, total, '監修・体制')

@register
def s_now(s, idx, total):
    s = slide(PAPER)
    kicker(s,'いま、起きていること')
    title(s,[[('AIは、一部の先進企業だけの話ではありません。',{'size':28,'color':INK,'bold':True})],
             [('あらゆる規模・業種の企業が、',{'size':28,'color':INK,'bold':True}),('もう日常業務で使っています。',{'size':28,'color':BLUE,'bold':True})]],line=1.28)
    textbox(s,MX,Inches(2.45),Inches(10.5),Inches(0.9),
        [[('「AIで何ができるか」を考える段階は終わり、「どの業務に、どう使い、どう定着させるか」を',{'size':15,'color':SUB})],
         [('競い合う段階に入りました。動き出した企業から、着実に差が開き始めています。',{'size':15,'color':SUB})]],line=1.7)
    inds=['建設・不動産','製造','運輸・物流','商社・小売','医療・介護','金融・保険','士業・コンサル','IT・通信','サービス・ホテル','官公庁・自治体']
    x=MX; y=Inches(3.7); maxx=EW-MX
    for p in inds:
        w=Inches(0.3+len(p)*0.155)
        if x+w>maxx: x=MX; y+=Inches(0.55)
        rect(s,x,y,w,Inches(0.42),fill=WHITE,line_color=LINE,rounded=True,radius=0.5)
        textbox(s,x,y,w,Inches(0.42),[[(p,{'size':12,'color':INK})]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        x+=w+Inches(0.16)
    footer(s, idx, total, '現状認識')

@register
def s_problem(s, idx, total):
    s = slide(PAPER)
    kicker(s,'多くの企業がぶつかる壁')
    title(s,[[('本当の課題は「ツールを入れること」ではなく、',{'size':27,'color':INK,'bold':True})],
             [('現場で',{'size':27,'color':INK,'bold':True}),('“使われ続けるか”＝定着',{'size':27,'color':BLUE,'bold':True}),('です。',{'size':27,'color':INK,'bold':True})]],line=1.26)
    items=[('01','各自バラバラに触っているだけ','個人で使えても、業務に落とし込めず、社内で活用が揃わない。'),
           ('02','導入しても、一部の人しか使わない','ツールを配ることと、現場で使われ続けることは別問題。'),
           ('03','研修しても、元の業務に戻る','「なるほど」で終わり、翌週には以前のやり方に戻ってしまう。')]
    n=3; gap=Inches(0.3); cw=(CW-gap*2)/3; y=Inches(2.7); h=Inches(2.5)
    for i,(no,t,d) in enumerate(items):
        x=MX+i*(cw+gap)
        rect(s,x,y,cw,h,fill=PANEL,line_color=LINE,rounded=True,radius=0.06)
        textbox(s,x+Inches(0.3),y+Inches(0.3),Inches(1),Inches(0.5),[[(no,{'size':22,'color':RGBColor(0xC2,0xCB,0xDA),'bold':True,'latin':'Arial'})]])
        textbox(s,x+Inches(0.3),y+Inches(0.95),cw-Inches(0.6),Inches(0.8),[[(t,{'size':15.5,'color':INK,'bold':True})]],line=1.3)
        textbox(s,x+Inches(0.3),y+Inches(1.65),cw-Inches(0.6),Inches(0.8),[[(d,{'size':12.5,'color':SUB})]],line=1.6)
    footer(s, idx, total, '課題')

@register
def s_fear(s, idx, total):
    s = slide(NAVY)
    kicker(s,'投資判断をためらう、最大の理由',dark=True)
    textbox(s,MX,Inches(2.2),Inches(6.4),Inches(2.4),
        [[('研修しても、',{'size':36,'color':WHITE,'bold':True})],
         [('結局、',{'size':36,'color':WHITE,'bold':True}),('誰も使わない。',{'size':36,'color':CREAM,'bold':True})]],line=1.45)
    textbox(s,EW-MX-Inches(5.0),Inches(2.35),Inches(5.0),Inches(2.6),
        [[('AI導入をためらう最大の理由は、費用でも難しさでもなく',{'size':15,'color':CREAM})],
         [('「お金をかけても、現場が元に戻る」という不安です。',{'size':15,'color':CREAM})],
         [('',{'size':8,'color':CREAM})],
         [('私たちが最初に解くのも、この“定着しない”問題。',{'size':15,'color':WHITE,'bold':True})],
         [('だから、研修の設計思想そのものが違います。',{'size':15,'color':WHITE,'bold':True})]],line=1.8)
    footer(s, idx, total, '本当の課題', dark=True)

@register
def s_value(s, idx, total):
    s = slide(NAVY)
    kicker(s,'私たちが売っているもの',dark=True)
    textbox(s,MX,Inches(2.3),Inches(6.6),Inches(2.4),
        [[('売るのは“研修”ではなく、',{'size':30,'color':WHITE,'bold':True})],
         [('“業務が変わった状態”。',{'size':34,'color':CREAM,'bold':True})]],line=1.45)
    textbox(s,EW-MX-Inches(4.9),Inches(2.4),Inches(4.9),Inches(2.6),
        [[('「使い方を教える研修」なら世の中に数多くあります。',{'size':15,'color':CREAM})],
         [('私たちが大切にするのは、研修後に',{'size':15,'color':CREAM})],
         [('“あの業務が実際に変わっている”状態まで',{'size':15,'color':WHITE,'bold':True})],
         [('持っていくこと。題材も御社の実業務そのものを使います。',{'size':15,'color':CREAM})]],line=1.8)
    footer(s, idx, total, '提供価値', dark=True)

@register
def s_compare(s, idx, total):
    s = slide(PAPER)
    kicker(s,'一般的なAI研修との違い')
    title(s,[[('「知っている」で終わらせず、',{'size':29,'color':INK,'bold':True}),('「使えている」まで。',{'size':29,'color':BLUE,'bold':True})]])
    rows=[('形式','座学・講義型のインプット／e-Learning','実務課題をその場で解く実践ワークショップ'),
          ('教材','汎用サンプル','御社の実業務（マスキング）を題材化'),
          ('成果物','ノート・受講証明書','参加者の数だけ“業務改善AI”が完成'),
          ('研修後','時間が経つと元の業務に戻る','3ヶ月伴走し、社内でPDCAを自走')]
    y=Inches(2.1); rh=Inches(0.86); c0=Inches(1.5); c1=Inches(4.7)
    x1=MX+c0; x2=x1+c1
    # headers
    textbox(s,x1,y,c1,Inches(0.4),[[('一般的なAI研修',{'size':12.5,'color':SUB,'bold':True})]])
    textbox(s,x2+Inches(0.25),y,c1+Inches(0.4),Inches(0.4),[[('弊社の実践型研修',{'size':12.5,'color':BLUE_D,'bold':True})]])
    hline(s,MX,y+Inches(0.42),CW,color=LINE,weight=1.2)
    yy=y+Inches(0.5)
    # highlight column bg
    rect(s,x2,yy,c1+Inches(0.55),rh*4,fill=BLUESOFT,rounded=True,radius=0.04)
    for lab,a,b in rows:
        textbox(s,MX,yy+Inches(0.16),c0,Inches(0.5),[[(lab,{'size':12,'color':FAINT})]])
        textbox(s,x1,yy+Inches(0.14),c1-Inches(0.2),Inches(0.6),[[(a,{'size':13.5,'color':SUB})]],line=1.4)
        textbox(s,x2+Inches(0.25),yy+Inches(0.14),c1+Inches(0.1),Inches(0.6),[[(b,{'size':13.5,'color':BLUE_D,'bold':True})]],line=1.4)
        yy+=rh
        if lab!='研修後': hline(s,MX,yy,CW,color=LINE2,weight=1.0)
    footer(s, idx, total, '一般研修との違い')

@register
def s_pillars(s, idx, total):
    s = slide(PAPER)
    kicker(s,'弊社が選ばれる3つの理由')
    title(s,[[('「定着」を、設計に組み込んでいます。',{'size':29,'color':INK,'bold':True})]])
    items=[('1','定着まで設計','座学で終わらせず、受講者の実業務をその場で題材に。研修後の運用まで伴走し、“使われ続ける状態”を作ります。'),
           ('2','事業を作ってきた人間が設計','研修専業ではなく、実際に事業を立ち上げてきたメンバーが、現場業務とROIの視点で設計します。'),
           ('3','効果の見える化','受講前後で「この業務が何分→何分」を測定。経営に報告できる形で成果を残します。')]
    n=3; gap=Inches(0.3); cw=(CW-gap*2)/3; y=Inches(2.35); h=Inches(2.9)
    for i,(no,t,d) in enumerate(items):
        x=MX+i*(cw+gap)
        rect(s,x,y,cw,h,fill=WHITE,line_color=LINE,rounded=True,radius=0.06)
        rect(s,x+Inches(0.32),y+Inches(0.34),Inches(0.5),Inches(0.5),fill=BLUE_D,rounded=True,radius=0.22)
        textbox(s,x+Inches(0.32),y+Inches(0.34),Inches(0.5),Inches(0.5),[[(no,{'size':16,'color':WHITE,'bold':True,'latin':'Arial'})]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        textbox(s,x+Inches(0.32),y+Inches(1.05),cw-Inches(0.64),Inches(0.8),[[(t,{'size':17,'color':INK,'bold':True})]],line=1.3)
        textbox(s,x+Inches(0.32),y+Inches(1.78),cw-Inches(0.6),Inches(1.0),[[(d,{'size':12.5,'color':SUB})]],line=1.65)
    footer(s, idx, total, '選ばれる理由')

@register
def s_retention(s, idx, total):
    s = slide(PAPER)
    kicker(s,'「また使われないのでは」への答え')
    title(s,[[('受講者が、',{'size':29,'color':INK,'bold':True}),('自分でAIを作って回し始める',{'size':29,'color':BLUE,'bold':True}),('まで。',{'size':29,'color':INK,'bold':True})]])
    textbox(s,MX,Inches(2.05),Inches(6.5),Inches(2.6),
        [[('「どうせ定着しない」は、皆さん最初に必ず言われます。',{'size':15,'color':SUB})],
         [('私たちが最も力を入れているのが、この“その後”です。',{'size':15,'color':SUB})]],line=1.7)
    bullets=[('受講者が自らGPTs等を作成し、研修後も“自走”する状態まで設計。'),
             ('配って終わりにせず、部署横断で日常的に使われる状態へ。'),
             ('「この業務が何分→何分」を測り、成果を組織に残す。')]
    by=Inches(3.05)
    for b in bullets:
        rect(s,MX,by+Pt(3),Inches(0.16),Inches(0.16),fill=BLUE,rounded=True,radius=0.3)
        textbox(s,MX+Inches(0.34),by,Inches(6.0),Inches(0.6),[[(b,{'size':13.5,'color':INK})]],line=1.45)
        by+=Inches(0.62)
    # right panel
    cx=EW-MX-Inches(4.4)
    rect(s,cx,Inches(2.0),Inches(4.4),Inches(3.0),fill=NAVY,rounded=True,radius=0.05)
    textbox(s,cx+Inches(0.4),Inches(2.45),Inches(3.6),Inches(1.0),
        [[('配って終わり、',{'size':22,'color':WHITE,'bold':True})],[('ではありません。',{'size':22,'color':WHITE,'bold':True})]],line=1.35)
    textbox(s,cx+Inches(0.4),Inches(3.75),Inches(3.65),Inches(1.1),
        [[('大企業でも、最初は“一部しか使わない”状態から、',{'size':12.5,'color':CREAM})],
         [('全部署が日常的に使う状態まで設計できます。',{'size':12.5,'color':CREAM})]],line=1.6)
    footer(s, idx, total, '定着の証明')

@register
def s_cases(s, idx, total):
    s = slide(PAPER)
    kicker(s,'導入企業の成果（実測）')
    title(s,[[('“浮いた時間”が、投資を回収します。',{'size':29,'color':INK,'bold':True})]])
    cases=[('建設業','安全書類の作成','25h','10h','月15時間削減・60%減','元請けごとに異なる書類をAIが8割作成。議事録も10h→3hに。'),
           ('会計事務所','顧客への月次レポート','24h','12h','月12時間削減・50%減','データ統合〜文面作成を定型化。海外向け提案資料も15h→5h。'),
           ('ホテル','シフト作成業務','10h','5h','月5時間削減・50%減','属人化していた作成を仕組み化。メニュー案作成も10h→6h。')]
    n=3; gap=Inches(0.3); cw=(CW-gap*2)/3; y=Inches(2.2); h=Inches(2.95)
    for i,(co,biz,a,b,delta,desc) in enumerate(cases):
        x=MX+i*(cw+gap)
        rect(s,x,y,cw,h,fill=WHITE,line_color=LINE,rounded=True,radius=0.06)
        textbox(s,x+Inches(0.3),y+Inches(0.28),cw-Inches(0.6),Inches(0.3),[[(co,{'size':12.5,'color':BLUE,'bold':True})]])
        textbox(s,x+Inches(0.3),y+Inches(0.6),cw-Inches(0.6),Inches(0.3),[[(biz,{'size':11.5,'color':FAINT})]])
        textbox(s,x+Inches(0.3),y+Inches(1.0),cw-Inches(0.6),Inches(0.7),
            [[(a,{'size':19,'color':FAINT,'latin':'Arial'}),('  →  ',{'size':15,'color':BLUE}),(b,{'size':34,'color':BLUE_D,'bold':True,'latin':'Arial'})]],anchor=MSO_ANCHOR.MIDDLE)
        rect(s,x+Inches(0.3),y+Inches(1.78),Inches(2.1),Inches(0.34),fill=RGBColor(0xF6,0xEC,0xE6),rounded=True,radius=0.5)
        textbox(s,x+Inches(0.3),y+Inches(1.78),Inches(2.1),Inches(0.34),[[(delta,{'size':11,'color':DELTA,'bold':True})]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        textbox(s,x+Inches(0.3),y+Inches(2.25),cw-Inches(0.6),Inches(0.7),[[(desc,{'size':11.5,'color':SUB})]],line=1.55)
    textbox(s,MX,Inches(5.45),CW,Inches(0.4),[[('※ いずれも受講後アンケート・実測値。業種・規模を問わず成果が出ています。',{'size':11,'color':SUB})]])
    footer(s, idx, total, '導入事例')

@register
def s_curriculum(s, idx, total):
    s = slide(PAPER)
    kicker(s,'研修カリキュラム')
    title(s,[[('全12時間（2時間×6回）。',{'size':28,'color':INK,'bold':True}),('すべて実務に直結。',{'size':28,'color':BLUE,'bold':True})]])
    textbox(s,MX,Inches(1.75),CW,Inches(0.4),[[('講義時間の多くを演習・ワークに充てるハンズオン構成。各回で“その場で動くもの”を作ります。',{'size':13,'color':SUB})]])
    data=[('1','生成AIの基礎とプロンプト設計','出力比較／業務プロンプト作成'),
          ('2','表計算 × AIによる業務効率化','コードによる自動化／データ整形・転記'),
          ('3','データ分析 & レポート・資料作成','レポート作成／スライド生成AI'),
          ('4','画像生成・音声入力 × AI活用','業務用ラフ案生成／短尺PR動画'),
          ('5','GPTs構築・ナレッジ活用','社内資料のQ&A化／RAG型質問応答'),
          ('6','全体振り返りと PoC プランニング','自社業務への適用設計／PoC決定')]
    y=Inches(2.35); rh=Inches(0.62)
    textbox(s,MX+Inches(0.0),y,Inches(0.6),Inches(0.3),[[('回',{'size':11,'color':FAINT})]])
    textbox(s,MX+Inches(0.7),y,Inches(4),Inches(0.3),[[('テーマ',{'size':11,'color':FAINT})]])
    textbox(s,MX+Inches(5.6),y,Inches(5),Inches(0.3),[[('主なワーク（実技）',{'size':11,'color':FAINT})]])
    hline(s,MX,y+Inches(0.32),CW,color=LINE,weight=1.2)
    yy=y+Inches(0.42)
    for no,t,w in data:
        textbox(s,MX,yy+Inches(0.1),Inches(0.6),Inches(0.4),[[(no,{'size':14,'color':BLUE,'bold':True,'latin':'Arial'})]])
        textbox(s,MX+Inches(0.7),yy+Inches(0.1),Inches(4.9),Inches(0.5),[[(t,{'size':13.5,'color':INK,'bold':True})]],line=1.2)
        textbox(s,MX+Inches(5.6),yy+Inches(0.1),Inches(5.3),Inches(0.5),[[(w,{'size':12.5,'color':SUB})]],line=1.2)
        yy+=rh
        hline(s,MX,yy,CW,color=LINE2,weight=1.0)
    footer(s, idx, total, '研修の中身')

@register
def s_process(s, idx, total):
    s = slide(PAPER)
    kicker(s,'導入の進め方')
    title(s,[[('お渡しして終わりにしない、',{'size':28,'color':INK,'bold':True}),('約半年の伴走。',{'size':28,'color':BLUE,'bold':True})]])
    steps=[('1','事前ヒアリング','業務・課題を把握し、効く領域を見立てる'),
           ('2','専用カリキュラム作成','御社の実業務を題材に内容を設計'),
           ('3','研修の実施','全12時間。オンライン中心、対面も可'),
           ('4','3ヶ月の定着フォロー','運用に伴走し、効果を測定・見える化')]
    n=4; gap=Inches(0.0); cw=CW/4; y=Inches(2.7)
    # connecting line
    hline(s,MX+Inches(0.4),y+Inches(0.25),CW-Inches(0.9),color=LINE,weight=1.4)
    for i,(no,t,d) in enumerate(steps):
        x=MX+i*cw
        rect(s,x,y,Inches(0.5),Inches(0.5),fill=BLUE_D,rounded=True,radius=0.5)
        textbox(s,x,y,Inches(0.5),Inches(0.5),[[(no,{'size':15,'color':WHITE,'bold':True,'latin':'Arial'})]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        textbox(s,x,y+Inches(0.75),cw-Inches(0.3),Inches(0.5),[[(t,{'size':15,'color':INK,'bold':True})]],line=1.25)
        textbox(s,x,y+Inches(1.25),cw-Inches(0.4),Inches(0.9),[[(d,{'size':12,'color':SUB})]],line=1.5)
    textbox(s,MX,Inches(5.1),CW,Inches(0.5),
        [[('最低5名から（雇用保険加入者）。',{'size':14,'color':INK,'bold':True}),
          ('  経営者・役員1名は無料でご招待します（5名以上のお申込み時）。',{'size':14,'color':SUB})]])
    footer(s, idx, total, '進め方')

@register
def s_subsidy(s, idx, total):
    s = slide(PAPER)
    kicker(s,'助成金の活用')
    title(s,[[('国の助成金で、研修費の負担を',{'size':28,'color':INK,'bold':True}),('大きく抑えられます。',{'size':28,'color':BLUE,'bold':True})]])
    textbox(s,MX,Inches(1.8),Inches(6.6),Inches(1.0),
        [[('人材開発支援助成金「事業展開等リスキリング支援コース」。',{'size':14,'color':SUB})],
         [('雇用保険の適用事業所が対象で、要件を満たした訓練が助成対象になります。',{'size':14,'color':SUB})]],line=1.6)
    # rate cards
    cards=[('中小企業','75','%','訓練経費の助成率'),('大企業','60','%','訓練経費の助成率')]
    for i,(seg,num,u,lab) in enumerate(cards):
        x=MX+i*Inches(3.35); y=Inches(3.05)
        rect(s,x,y,Inches(3.1),Inches(1.7),fill=PANEL,line_color=LINE,rounded=True,radius=0.06)
        textbox(s,x+Inches(0.3),y+Inches(0.22),Inches(2.5),Inches(0.3),[[(seg,{'size':12,'color':SUB,'bold':True})]])
        textbox(s,x+Inches(0.28),y+Inches(0.5),Inches(2.6),Inches(0.8),[[('最大 ',{'size':14,'color':INK}),(num,{'size':38,'color':BLUE_D,'bold':True,'latin':'Arial'}),(u,{'size':20,'color':BLUE,'bold':True})]],anchor=MSO_ANCHOR.MIDDLE)
        textbox(s,x+Inches(0.3),y+Inches(1.28),Inches(2.6),Inches(0.3),[[(lab,{'size':11.5,'color':FAINT})]])
    # right panel
    cx=EW-MX-Inches(4.3)
    rect(s,cx,Inches(1.8),Inches(4.3),Inches(2.95),fill=BLUESOFT,line_color=RGBColor(0xDD,0xE6,0xF6),rounded=True,radius=0.05)
    pts=['訓練経費に加え、訓練中の賃金も助成対象（時間に応じ加算）。',
         '申請は実績豊富な提携社労士法人が代行。御社の手間は最小限に。',
         '弊社は社労士の紹介料を一切いただきません。']
    py=Inches(2.1)
    for p in pts:
        rect(s,cx+Inches(0.32),py+Pt(3),Inches(0.14),Inches(0.14),fill=BLUE,rounded=True,radius=0.3)
        textbox(s,cx+Inches(0.6),py,Inches(3.4),Inches(0.7),[[(p,{'size':12,'color':BLUE_D})]],line=1.5)
        py+=Inches(0.78)
    textbox(s,MX,Inches(5.05),CW,Inches(0.5),
        [[('※ 助成額・助成率は企業区分・年度・申請内容により異なり、労働局の審査・支給決定が必要です。受給を保証するものではありません。',{'size':10.5,'color':FAINT})]],line=1.4)
    footer(s, idx, total, '助成金')

@register
def s_price(s, idx, total):
    s = slide(PAPER)
    kicker(s,'費用')
    title(s,[[('受講料と、助成金活用時の負担イメージ。',{'size':28,'color':INK,'bold':True})]])
    # left: standard price
    y=Inches(2.3); h=Inches(2.0)
    rect(s,MX,y,Inches(5.2),h,fill=WHITE,line_color=LINE,rounded=True,radius=0.06)
    textbox(s,MX+Inches(0.4),y+Inches(0.3),Inches(4),Inches(0.3),[[('標準受講料',{'size':12.5,'color':SUB})]])
    textbox(s,MX+Inches(0.38),y+Inches(0.65),Inches(4.6),Inches(0.8),[[('¥300,300',{'size':38,'color':INK,'bold':True,'latin':'Arial'}),(' / 人（税込）',{'size':14,'color':SUB})]],anchor=MSO_ANCHOR.MIDDLE)
    textbox(s,MX+Inches(0.4),y+Inches(1.45),Inches(4.6),Inches(0.5),[[('全12時間（2h×6回）＋ 研修後3ヶ月のフォローアップ込み。最低5名から。',{'size':11.5,'color':SUB})]],line=1.4)
    # right: subsidized estimate
    cx=MX+Inches(5.55)
    rect(s,cx,y,Inches(5.45),h,fill=NAVY,rounded=True,radius=0.06)
    textbox(s,cx+Inches(0.4),y+Inches(0.3),Inches(4.6),Inches(0.4),[[('助成金を活用した場合の自己負担イメージ（中小企業）',{'size':12,'color':CREAM})]],line=1.2)
    textbox(s,cx+Inches(0.38),y+Inches(0.72),Inches(4.8),Inches(0.8),[[('約 ¥98,675',{'size':38,'color':WHITE,'bold':True,'latin':'Arial'}),(' / 人',{'size':14,'color':CREAM})]],anchor=MSO_ANCHOR.MIDDLE)
    textbox(s,cx+Inches(0.4),y+Inches(1.5),Inches(4.7),Inches(0.4),[[('社労士費用込みの目安（要件を満たし支給決定を受けた場合）。経営者1名は無料招待。',{'size':10.5,'color':CREAM})]],line=1.35)
    textbox(s,MX,Inches(4.55),CW,Inches(0.6),
        [[('浮いた時間で回収できる投資です。',{'size':15,'color':INK,'bold':True}),
          ('  月15時間かかっていた業務が半分になれば、年間で約90時間を本来の業務に戻せます。',{'size':14,'color':SUB})]],line=1.5)
    textbox(s,MX,Inches(5.25),CW,Inches(0.4),[[('※ 上記は一例の目安です。受給可否・金額は要件充足と労働局の審査・支給決定によります。',{'size':10.5,'color':FAINT})]])
    footer(s, idx, total, '費用')

@register
def s_whynow(s, idx, total):
    s = slide(PAPER)
    kicker(s,'なぜ“今”か')
    title(s,[[('差がつくのは、ツールの新しさではなく',{'size':28,'color':INK,'bold':True})],
             [('“',{'size':28,'color':INK,'bold':True}),('土台を作った時期',{'size':28,'color':BLUE,'bold':True}),('”です。',{'size':28,'color':INK,'bold':True})]],line=1.25)
    items=[('AIの進化は速い','ツールは日進月歩。差がつくのは「社員が日常業務で使いこなせるか」という活用力です。'),
           ('複利で開く差','AIは走りながら育てるもの。早く始めた企業ほど、社内にノウハウが蓄積していきます。'),
           ('活用の好機に期限','助成金の活用機会には期限があり、準備にも時間を要します。検討は早いほど有利です。')]
    n=3; gap=Inches(0.3); cw=(CW-gap*2)/3; y=Inches(2.6); h=Inches(2.3)
    for i,(t,d) in enumerate(items):
        x=MX+i*(cw+gap)
        rect(s,x,y,cw,h,fill=WHITE,line_color=LINE,rounded=True,radius=0.06)
        rect(s,x+Inches(0.3),y+Inches(0.34),Inches(0.34),Pt(3),fill=GOLD)
        textbox(s,x+Inches(0.3),y+Inches(0.5),cw-Inches(0.6),Inches(0.5),[[(t,{'size':16,'color':BLUE_D,'bold':True})]])
        textbox(s,x+Inches(0.3),y+Inches(1.1),cw-Inches(0.6),Inches(1.1),[[(d,{'size':12.5,'color':SUB})]],line=1.65)
    footer(s, idx, total, '今、始める理由')

@register
def s_summary(s, idx, total):
    s = slide(NAVY)
    kicker(s,'まとめ',dark=True)
    textbox(s,MX,Inches(2.1),Inches(7.0),Inches(2.6),
        [[('完璧を待たず、まず始める。',{'size':32,'color':WHITE,'bold':True})],
         [('最高のスタートを、今。',{'size':32,'color':CREAM,'bold':True})]],line=1.45)
    pts=['研修して終わりにせず、“現場が変わる”まで設計','受講前後で効果を測り、経営に見える形で残す','助成金の活用で、投資負担を抑えて始められる']
    py=Inches(4.5)
    for p in pts:
        rect(s,MX,py+Pt(3),Inches(0.16),Inches(0.16),fill=CREAM,rounded=True,radius=0.3)
        textbox(s,MX+Inches(0.34),py,Inches(10),Inches(0.4),[[(p,{'size':14,'color':WHITE})]])
        py+=Inches(0.5)
    footer(s, idx, total, 'まとめ', dark=True)

@register
def s_nextstep(s, idx, total):
    s = slide(PAPER)
    kicker(s,'次のステップ')
    title(s,[[('ご導入までの流れ',{'size':29,'color':INK,'bold':True})]])
    steps=[('01','無料相談・デモ','御社の業務に合わせた活用イメージをご覧いただきます。'),
           ('02','活用設計・効果試算','どの業務から始めると効くか、ROIを一緒に試算します。'),
           ('03','お申込み・助成金申請','提携社労士が申請を代行。計画申請から伴走します。'),
           ('04','研修開始・定着支援','専用カリキュラムで研修を実施し、3ヶ月伴走します。')]
    y=Inches(2.25); rh=Inches(0.92)
    for i,(no,t,d) in enumerate(steps):
        yy=y+i*rh
        rect(s,MX,yy,Inches(0.62),Inches(0.62),fill=BLUESOFT,rounded=True,radius=0.16)
        textbox(s,MX,yy,Inches(0.62),Inches(0.62),[[(no,{'size':15,'color':BLUE_D,'bold':True,'latin':'Arial'})]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        textbox(s,MX+Inches(0.95),yy+Inches(0.05),Inches(4),Inches(0.5),[[(t,{'size':16,'color':INK,'bold':True})]],anchor=MSO_ANCHOR.MIDDLE)
        textbox(s,MX+Inches(5.0),yy,Inches(6.0),Inches(0.62),[[(d,{'size':13,'color':SUB})]],anchor=MSO_ANCHOR.MIDDLE,line=1.4)
        if i<3: hline(s,MX+Inches(0.95),yy+rh-Inches(0.15),CW-Inches(0.95),color=LINE2,weight=1.0)
    footer(s, idx, total, '次のステップ')

@register
def s_contact(s, idx, total):
    s = slide(NAVY)
    logo(s, MX, Inches(0.9), dark=True)
    textbox(s,MX,Inches(2.2),Inches(10),Inches(0.8),[[('お気軽にご相談ください。',{'size':30,'color':WHITE,'bold':True})]])
    textbox(s,MX,Inches(3.05),Inches(9),Inches(0.8),
        [[('ご質問・無料デモのご希望など、御社に合わせた進め方をご提案します。',{'size':15,'color':CREAM})]],line=1.6)
    info=[('会社名','株式会社AIスキル'),('代表者','立石 亮介'),
          ('所在地','東京都渋谷区神南1-11-4 FPGリンクス神南 5階'),
          ('事業内容','AI活用サポート事業／教育事業／コミュニティ事業／SNS事業')]
    iy=Inches(4.1)
    for k,v in info:
        textbox(s,MX,iy,Inches(1.6),Inches(0.4),[[(k,{'size':11,'color':RGBColor(0x8F,0xA3,0xC6)})]])
        textbox(s,MX+Inches(1.7),iy,Inches(9),Inches(0.4),[[(v,{'size':13.5,'color':WHITE,'bold':True})]])
        iy+=Inches(0.5)
    footer(s, idx, total, 'お問い合わせ', dark=True)

# -------- appendix
@register
def s_appendix_div(s, idx, total):
    s = slide(PAPER)
    textbox(s,MX,Inches(2.9),Inches(4),Inches(0.4),[[('Appendix',{'size':12,'color':BLUE,'bold':True,'spacing':2})]])
    title(s,[[('付録：詳細資料',{'size':34,'color':INK,'bold':True})]],y=Inches(3.3))
    textbox(s,MX,Inches(4.2),Inches(9),Inches(0.5),[[('料金の内訳・よくあるご質問。商談に合わせて必要なページのみご参照ください。',{'size':14,'color':SUB})]])
    footer(s, idx, total, '付録')

@register
def s_price_detail(s, idx, total):
    s = slide(PAPER)
    kicker(s,'付録 ｜ 料金の内訳')
    title(s,[[('助成金活用時の自己負担イメージ（1人あたり・目安）',{'size':24,'color':INK,'bold':True})]])
    rows=[('標準受講料','¥300,300','¥300,300'),
          ('訓練経費の助成（75%／60%）','約 −¥225,225','約 −¥180,180'),
          ('賃金助成（訓練時間に応じ）','約 −¥12,000','約 −¥6,000'),
          ('社労士費用（助成額の約15%）','約 ＋¥35,600','約 ＋¥27,900'),
          ('自己負担イメージ','約 ¥98,675','約 ¥142,000')]
    y=Inches(2.0); rh=Inches(0.68); c0=Inches(5.4); c1=Inches(2.8)
    textbox(s,MX+c0,y,c1,Inches(0.3),[[('中小企業',{'size':12,'color':SUB,'bold':True})]],align=PP_ALIGN.RIGHT)
    textbox(s,MX+c0+c1,y,c1,Inches(0.3),[[('大企業',{'size':12,'color':BLUE_D,'bold':True})]],align=PP_ALIGN.RIGHT)
    hline(s,MX,y+Inches(0.34),CW,color=LINE,weight=1.2)
    yy=y+Inches(0.44)
    for i,(lab,a,b) in enumerate(rows):
        last=(i==len(rows)-1)
        if last: rect(s,MX-Inches(0.1),yy-Inches(0.04),CW+Inches(0.2),rh,fill=BLUESOFT,rounded=True,radius=0.1)
        textbox(s,MX,yy+Inches(0.13),c0,Inches(0.4),[[(lab,{'size':13,'color':INK if last else SUB,'bold':last})]])
        textbox(s,MX+c0,yy+Inches(0.1),c1,Inches(0.4),[[(a,{'size':14 if last else 13,'color':INK if last else SUB,'bold':last,'latin':'Arial'})]],align=PP_ALIGN.RIGHT)
        textbox(s,MX+c0+c1,yy+Inches(0.1),c1,Inches(0.4),[[(b,{'size':14 if last else 13,'color':BLUE_D if last else SUB,'bold':last,'latin':'Arial'})]],align=PP_ALIGN.RIGHT)
        yy+=rh
        if not last: hline(s,MX,yy,CW,color=LINE2,weight=1.0)
    textbox(s,MX,Inches(6.05),CW,Inches(0.7),[[('※ 経費助成は訓練時間10〜100時間で上限30万円/人。動画教材を用いた研修は賃金助成の対象外。金額・助成率は企業区分・年度・申請内容により変動し、労働局の審査・支給決定が必要です（受給を保証するものではありません）。',{'size':10,'color':FAINT})]],line=1.45)
    footer(s, idx, total, '付録｜料金')

@register
def s_faq(s, idx, total):
    s = slide(PAPER)
    kicker(s,'付録 ｜ よくあるご質問')
    title(s,[[('よくあるご質問',{'size':28,'color':INK,'bold':True})]])
    qa=[('オンラインですか、対面ですか？','最低5名以上で基本オンライン。10名以上なら対面も可能です（交通費別途）。アーカイブ視聴も可。'),
        ('研修期間はどのくらい？','全12時間（2h×6回が中心）。1ヶ月集中も2〜3ヶ月も、御社の状況に合わせて調整します。'),
        ('社員のITリテラシーが不安です。','幅広い年代で好評をいただいています。実業務を題材にするため「自分の仕事でどう使うか」がつかめます。'),
        ('個人情報は使って大丈夫？','実データはマスキングしたサンプルを使用。実業務に近い形で安全に演習できます。'),
        ('すでに他社研修を受けています。','内容が異なれば助成金は複数回の活用が可能な場合があります。一度内容をご確認ください。'),
        ('申請が難しそうです。','提携社労士が代行。対象判定から書類準備まで支援します（弊社の紹介料なし）。')]
    cols=2; gap=Inches(0.6); cw=(CW-gap)/2; y=Inches(2.0); rh=Inches(1.15)
    for i,(q,a) in enumerate(qa):
        c=i%2; r=i//2
        x=MX+c*(cw+gap); yy=y+r*rh
        textbox(s,x,yy,cw,Inches(0.4),[[('Q. ',{'size':13.5,'color':BLUE,'bold':True}),(q,{'size':13.5,'color':INK,'bold':True})]],line=1.25)
        textbox(s,x,yy+Inches(0.38),cw,Inches(0.7),[[(a,{'size':12,'color':SUB})]],line=1.5)
    footer(s, idx, total, '付録｜FAQ')

# ---------- build
total=len(SLIDES)
for i,fn in enumerate(SLIDES, start=1):
    fn(None, i, total)

out=os.path.join(os.path.dirname(os.path.abspath(__file__)), "ai-training-proposal.pptx")
prs.save(out)
print("saved", out, "slides:", total)
