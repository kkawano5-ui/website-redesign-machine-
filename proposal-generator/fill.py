# -*- coding: utf-8 -*-
"""テンプレPPTXの可変箇所を差し替える。
- テキスト：クラン向けの実文字列をアンカーにして置換（テーブルはスキップ）
- 診断テーブル：ヘッダ['項目','現状','評価','改善ポイント']の表を見つけて8行を上書き
デザイン（配色・レイアウト）はテンプレのまま維持される。"""

from pptx import Presentation

DIAG_HEADER = ["項目", "現状", "評価", "改善ポイント"]


def _set_paragraph_text(p, new_text):
    """段落テキストを差し替え（先頭ランの書式を維持し、残りランを空に）。"""
    if p.runs:
        p.runs[0].text = new_text
        for r in p.runs[1:]:
            r.text = ""
    else:
        run = p.add_run()
        run.text = new_text


def _replace_in_textframe(tf, subs):
    for p in tf.paragraphs:
        full = "".join(r.text for r in p.runs)
        if not full:
            continue
        new = full
        for old, rep in subs:
            if old in new:
                new = new.replace(old, rep)
        if new != full:
            _set_paragraph_text(p, new)


def _iter_text_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # group
            for s in _iter_text_shapes(sh.shapes):
                yield s
        elif getattr(sh, "has_table", False):
            continue  # テーブルは別処理（診断）／固定表（比較）は触らない
        elif sh.has_text_frame:
            yield sh


def _fill_diagnosis(prs, diag_rows):
    for slide in prs.slides:
        for sh in slide.shapes:
            if not getattr(sh, "has_table", False):
                continue
            t = sh.table
            header = [t.cell(0, c).text.strip() for c in range(len(t.columns))]
            if header[:4] != DIAG_HEADER:
                continue
            for i, (item, cur, ev, imp) in enumerate(diag_rows, start=1):
                if i >= len(t.rows):
                    break
                for col, val in ((1, cur), (2, ev), (3, imp)):
                    cell = t.cell(i, col)
                    ps = cell.text_frame.paragraphs
                    _set_paragraph_text(ps[0], val)
                    for extra in ps[1:]:
                        _set_paragraph_text(extra, "")
            return True
    return False


def build_subs(data):
    """クラン・テンプレのアンカー文字列 → 新しい店舗の値。"""
    chips = data["chips"]          # 6個
    themes_joined = "・".join(data["themes"])
    return [
        # 検索KWチップ（6個・アンカー＝クランの実文字列）
        ("近くの不動産会社", chips[0]),
        ("鶴見 不動産売却", chips[1]),
        ("鶴見 マンション売却", chips[2]),
        ("相続不動産 相談", chips[3]),
        ("横浜 不動産買取", chips[4]),
        ("借地権 相談 横浜", chips[5]),
        # タイトル横の例文2行目
        ("横浜 相続不動産", data["line2"]),
        # テーマ語（スライド8・9に登場）
        ("売却・買取・相続・リノベ・借地権", themes_joined),
        # エリア表記（KW例ヘッダ・期待効果）
        ("鶴見・横浜・川崎", data["area_label"]),
        # 会社名（タイトル・診断見出し・推奨見出し）
        ("クラン株式会社", data["company"]),
        # 日付
        ("2026年5月", data["date"]),
    ]


def fill(template_path, out_path, data):
    prs = Presentation(template_path)
    subs = build_subs(data)
    for slide in prs.slides:
        for sh in _iter_text_shapes(slide.shapes):
            _replace_in_textframe(sh.text_frame, subs)
    ok = _fill_diagnosis(prs, data["diag_rows"])
    prs.save(out_path)
    return ok
